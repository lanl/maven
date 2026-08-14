
'''
###############################################################################
# This program is Open-Source under the BSD-3 License.                        #
#                                                                             #
# © 2026. Triad National Security, LLC. All rights reserved. O5152            #
#                                                                             #
# This program was produced under U.S. Government contract 89233218CNA000001  #
# for Los Alamos National Laboratory (LANL), which is operated by Triad       #
# National Security, LLC for the U.S. Department of Energy/National Nuclear   #
# Security Administration. All rights in the program are reserved by Triad    #
# National Security, LLC, and the U.S. Department of Energy/National Nuclear  #
# Security Administration. The Government is granted for itself and others    #
# acting on its behalf a nonexclusive, paid-up, irrevocable worldwide license #
# in this material to reproduce, prepare. derivative works, distribute copies #
# to the public, perform publicly and display publicly, and to permit others  #
# to do so.                                                                   #
# Redistribution and use in source and binary forms, with or without          #
# modification, are permitted provided that the following conditions are met: #
#                                                                             #  
# Redistributions of source code must retain the above copyright notice, this #
# list of conditions and the following disclaimer.                            #
# Redistributions in binary form must reproduce the above copyright notice,   #
# this list of conditions and the following disclaimer in the documentation   #
# and/or other materials provided with the distribution.                      #
# Neither the name of the copyright holder nor the names of its contributors  #
# may be used to endorse or promote products derived from this software       #
# without specific prior written permission.                                  #
#                                                                             #
# THIS SOFTWARE IS PROVIDED BY THE COPYRIGHT HOLDERS AND CONTRIBUTORS "AS IS" #
# AND ANY EXPRESS OR IMPLIED WARRANTIES, INCLUDING, BUT NOT LIMITED TO, THE   #
# IMPLIED WARRANTIES OF MERCHANTABILITY AND FITNESS FOR A PARTICULAR PURPOSE  #
# ARE DISCLAIMED. IN NO EVENT SHALL THE COPYRIGHT HOLDER OR CONTRIBUTORS BE   #
# LIABLE FOR ANY DIRECT, INDIRECT, INCIDENTAL, SPECIAL, EXEMPLARY, OR         #
# CONSEQUENTIAL DAMAGES (INCLUDING, BUT NOT LIMITED TO, PROCUREMENT OF        #
# SUBSTITUTE GOODS OR SERVICES; LOSS OF USE, DATA, OR PROFITS; OR BUSINESS    #
# INTERRUPTION) HOWEVER CAUSED AND ON ANY THEORY OF LIABILITY, WHETHER IN     #
# CONTRACT, STRICT LIABILITY, OR TORT (INCLUDING NEGLIGENCE OR OTHERWISE)     #
# ARISING IN ANY WAY OUT OF THE USE OF THIS SOFTWARE, EVEN IF ADVISED OF      #
# THE POSSIBILITY OF SUCH DAMAGE.                                             #
###############################################################################


User provides required inputs and an LLM parses the responses to populate the required fields for:
1) MAVEN Datasheet
2) Findability Metadata
3) AI-Ready Data
4) DSI Move


Authors:
Christopher W Johnson (cwj@lanl.gov)
Vedant Iyer (iyer@lanl.gov)
'''

import json
import yaml
from typing import Any, Dict, List, Tuple
from pathlib import Path
import os
import streamlit as st
import pypdf
import docx
from pptx import Presentation
import pandas as pd
import shutil
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle
import html
import subprocess
import shlex
from pathlib import PurePosixPath
import getpass
import socket
import re
import tempfile
from openai import OpenAI
import httpx
# from linkml_runtime import SchemaView
from datetime import UTC, datetime
import requests
import urllib3
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

from dsi.dsi import DSI
from dsi.sync import Sync

from maven.ursa_autofill import (
    AGENT_META_COLUMNS,
    build_followup_questions,
    load_agent_meta,
    merge_autofill_result,
    run_followup_autofill,
    run_initial_autofill,
    summarize_autofill,
)

from maven.tier1_genesis_dc_agent import run_tier1_catalog
from ursa.agents.chat_agent import ChatAgent
try:
    from ursa.util.http import inject_truststore_into_ssl
    inject_truststore_into_ssl()
except ImportError:
    raise ImportError("Ensure you have ursa-ai>=0.15.8 downloaded from pypi")
from langchain.chat_models import init_chat_model

curr_dir = Path(__file__).parent
files_dir = curr_dir / "files_template"

# GENESIS v1.2 datacard files
GENESIS_MISSION_DATA_CARD_BREAKDOWN = files_dir / "genesis_dc_breakdown.yaml"
GENESIS_MISSION_DATA_CARD_YAML = files_dir / "genesis_dc_v1.2_YAML.yaml"
GENESIS_MISSION_DATA_CARD_MD = files_dir / "genesis_dc_v1.2_MD.md"
GENESIS_MISSION_DATA_CARD_REFERENCE = files_dir / "genesis_dc_v1.2_reference_guide.md"

# Genesis datasheet 
datasheet_file = files_dir / "datasheet_sections.yaml"
SECTIONS = yaml.safe_load(datasheet_file.read_text(encoding="utf-8"))

REVIEW_SECTION_IDX = -1
FOLLOWUP_SECTION_IDX = -2

SECTION_BY_IDX = {s["section_idx"]: s for s in SECTIONS}
ALL_SECTION_IDXS = [s["section_idx"] for s in SECTIONS]
ALL_QUESTIONS = [q for s in SECTIONS for q in s["questions"]]


MASTER_DB_PATH = "maven.db"
PROJECTS_TABLE = "projects"

DATASHEET_TABLE = "datasheet"
FILE_POINTERS_TABLE = "metadata_file_paths"

TIER_2_TABLE = "data_and_hpc_info"

CONFIG_DIR = Path.home() / ".maven_config"
CONFIG_FILE = CONFIG_DIR / "files_location.txt"
API_KEYS_FILE = CONFIG_DIR / "ai_api_keys.txt"

MAVEN_FOLDER = "maven_files"

remote_script = curr_dir / "remote_move.py"
REMOTE_MOVE_SCRIPT = remote_script.read_text(encoding="utf-8")

remote_endpoint_script = curr_dir / "remote_register_endpoint.py"
REMOTE_REGISTER_ENDPOINT_SCRIPT = remote_endpoint_script.read_text(encoding="utf-8")

is_remote = any(os.environ.get(x) for x in ["SSH_CONNECTION", "SSH_CLIENT", "SSH_TTY"])

TEMP = 0.2

def get_maven_dir() -> Path | None:
    if not CONFIG_FILE.exists():
        return None

    path_str = CONFIG_FILE.read_text(encoding="utf-8").strip()
    if not path_str:
        return None

    path = Path(path_str).expanduser()

    if not path.is_dir() or not os.access(path, os.R_OK | os.W_OK | os.X_OK):
        return None

    return path


def configure_chat_agent() -> ChatAgent:
    # Configure model
    llm = init_chat_model(model=os.getenv("AI_MODEL"),
                          base_url=os.getenv("AI_API_URL"),
                          api_key=os.getenv("AI_API_KEY"),
                          temperature=TEMP
                          )
    # Setup workspace and thread for conversation persistence
    workspace = Path(get_maven_dir()) / "ursa_workspace"
    workspace.mkdir(parents=True, exist_ok=True)
    # Create ChatAgent with conversation state
    chat_agent = ChatAgent(llm=llm, workspace=workspace, autosave_metrics=False)
    return chat_agent


def save_maven_dir(path_str: str) -> bool:
    path = Path(path_str).expanduser()

    if not path.is_dir():
        st.error("That path is not a valid directory or does not exist.")
        return False

    if not path.is_absolute():
        st.error("Please provide an absolute directory path.")
        return False

    if not os.access(path, os.R_OK | os.W_OK | os.X_OK):
        st.error("You do not have read/write/access permission for that directory.")
        return False

    if path.name != MAVEN_FOLDER:
        path = path / MAVEN_FOLDER
    path.mkdir(parents=True, exist_ok=True)

    old_dir = get_maven_dir()
    if old_dir is not None and old_dir.resolve() != path.resolve():
        for item in old_dir.iterdir():
            target = path / item.name
            if target.exists():
                st.warning(f"Skipping move of {item.name} as it already exists in {path}.")
                continue
            shutil.move(str(item), str(target))

    CONFIG_DIR.mkdir(parents=True, exist_ok=True)
    CONFIG_FILE.write_text(str(path.resolve()) + "\n")
    return True


def load_env_keys() -> None:
    if API_KEYS_FILE.exists():
        for line in API_KEYS_FILE.read_text(encoding="utf-8").splitlines():
            line = line.strip()

            if not line or "=" not in line:
                return False

            key, value = line.split("=", 1)
            os.environ[key.strip()] = value.strip()
        if not os.environ.get("AI_API_KEY") or not os.environ.get("AI_API_URL") or not os.environ.get("AI_MODEL"):
            return False
        return True
    return False


def get_db(db_path: str) -> DSI:
    return DSI(db_path, check_same_thread=False, silence_messages=True)


def all_question_columns() -> List[str]:
    cols: List[str] = []
    for section in SECTIONS:
        for q in section["questions"]:
            if q["id"] != "short_project_title":
                cols.append(q["id"])
    cols.extend(AGENT_META_COLUMNS)
    cols.append("context_files_text")
    cols.append("ROSY_ID")
    cols.append("ROSY_Z_NUMBER")
    return cols


def make_tier1_db_name(short_title: str, only_name: bool = False) -> str:
    cleaned = short_title.strip().lower().replace(" ", "_").replace("/", "_")
    if only_name:
        return f"{cleaned}_tier1.db"
    maven_dir = get_maven_dir()
    return str(maven_dir / f"{cleaned}_tier1.db")


def make_tier2_db_name(short_title: str, only_name: bool = False) -> str:
    cleaned = short_title.strip().lower().replace(" ", "_").replace("/", "_")
    if only_name:
        return f"{cleaned}_tier2.db"
    maven_dir = get_maven_dir()
    return str(maven_dir / f"{cleaned}_tier2.db")


def get_master_db_name() -> str:
    maven_dir = get_maven_dir()
    return str(maven_dir / MASTER_DB_PATH)


def create_master_db():
    store = get_db(get_master_db_name())
    master_dict = {"project_id": "", "project_name": "", "tier1_db_path": "",
                   "tier2_db_path": "", "has_moved": ""}
    store.read(master_dict, "Collection", table_name=PROJECTS_TABLE)
    store.close()


def create_tier1_db(db_path: str):
    question_cols = all_question_columns()
    project_dict = {x: [""] for x in question_cols}

    store = get_db(db_path)
    store.read(project_dict, "Collection", DATASHEET_TABLE)
    store.close()


def create_tier2_db(db_path: str):
    store = get_db(db_path)
    tier2_dict = {"local_data_path": "", "username": "", "hpc_system": "",
                  "hpc_staging_space": "", "hpc_campaign_space": "", "user_group": "", 
                  "access_permissions": "", "diana_endpoint": "", "contact_email": ""}
    store.read(tier2_dict, "Collection", table_name=TIER_2_TABLE)
    store.close()


def list_projects() -> pd.DataFrame:
    store = get_db(get_master_db_name())
    project_df = store.query(f"SELECT project_id, project_name FROM {PROJECTS_TABLE} ORDER BY project_id ASC", True)
    store.close()
    return project_df


def get_tier1_db_path(project_id: int, only_name: bool = False) -> str:
    store = get_db(get_master_db_name())
    df = store.query(f"SELECT tier1_db_path FROM {PROJECTS_TABLE} WHERE project_id = {project_id}", True)
    store.close()
    if only_name:
        return df.iloc[0, 0]
    maven_dir = get_maven_dir()
    return str(maven_dir / df.iloc[0, 0])


def get_tier1_YAML_MD_path(project_id: int, only_name: bool = False) -> str:
    store = get_db(get_master_db_name())
    df = store.query(f"SELECT tier1_db_path FROM {PROJECTS_TABLE} WHERE project_id = {project_id}", True)
    store.close()
    if only_name:
        return df.iloc[0, 0]
    maven_dir = get_maven_dir()
    return str(maven_dir / df.iloc[0, 0]).replace(".db", ".yaml")


def get_tier2_db_path(project_id: int, only_name: bool = False) -> str:
    store = get_db(get_master_db_name())
    df = store.query(f"SELECT tier2_db_path FROM {PROJECTS_TABLE} WHERE project_id = {project_id}", True)
    store.close()
    if only_name:
        return df.iloc[0, 0]
    maven_dir = get_maven_dir()
    return str(maven_dir / df.iloc[0, 0])


def get_datasheet(qid: int, df_return=False) -> pd.DataFrame | Dict[str, Any]:
    tier1_db = get_tier1_db_path(qid)

    store = get_db(tier1_db)
    df = store.get_table(DATASHEET_TABLE, True)
    store.close()

    if df_return:
        return df
    return df.iloc[0].to_dict()


def update_datasheet(qid: int, updates: Dict[str, Any]) -> None:
    if not updates:
        return

    tier1_db = get_tier1_db_path(qid)
    store = get_db(tier1_db)

    query = (f'UPDATE {DATASHEET_TABLE} SET ' + ", ".join(f'"{k}" = ?' for k in updates))
    store.query(query, params=list(updates.values()))
    store.close()


def get_tier1_table(project_id: int, update=False, check_exists=False) -> Dict[str, pd.DataFrame] | bool:
    tier1_db = get_tier1_db_path(project_id)
    store = get_db(tier1_db)

    curr_tables = store.list(True)
    curr_tables = [t for t in curr_tables if t not in [DATASHEET_TABLE, FILE_POINTERS_TABLE, "filesystem", "federated"]]
    if check_exists:
        store.close()
        if not curr_tables:
            return False
        return True

    all_tbls = {}
    for tbl in curr_tables:
        df = store.get_table(tbl, collection=True, update=update)
        all_tbls[tbl] = df
    store.close()
    return all_tbls


def update_tier1_table(project_id: int, update_tier1_df: pd.DataFrame):
    tier1_db = get_tier1_db_path(project_id)
    store = get_db(tier1_db)
    store.update(update_tier1_df)
    store.close()


def delete_project(qid: int) -> None:
    tier1_path = get_tier1_db_path(qid)
    if os.path.exists(tier1_path):
        os.remove(tier1_path)

    tier2_path = get_tier2_db_path(qid)
    if os.path.exists(tier2_path):
        os.remove(tier2_path)

    store = get_db(get_master_db_name())
    projects_df = store.get_table(PROJECTS_TABLE, True, True)

    if len(projects_df) == 1:
        store.close()
        os.remove(get_master_db_name())
    else:
        projects_df["project_id"] = projects_df["project_id"].astype(int)
        projects_df = projects_df[projects_df["project_id"] != qid]
        store.update(projects_df)
        store.close()


# def get_tier1_fields() -> Tuple[Dict[str, str], Dict[str, Dict[str, str]]]:
#     if not CARD_CLASS_BASED.is_file():
#         st.error("Tier 1 metadata fields file does not exist")
#         st.stop()

#     all_classes = {}
#     sv = SchemaView(CARD_CLASS_BASED)
#     for class_name, cls in sv.all_classes().items():
#         if cls.abstract or class_name.lower() == "anyvalue":
#             continue

#         columns = {}
#         required_columns = []
#         for slot in sv.class_induced_slots(class_name):
#             if slot.name.lower() == "anyvalue":
#                 continue
#             columns[slot.name] = {'description': slot.description, 'required': True if slot.required else False}
#             if slot.required:
#                 required_columns.append(slot.name)
#         class_dict = {"description": cls.description, "columns": columns, "required_columns": required_columns}
#         all_classes[class_name] = class_dict

#     slots_dict = {slot_name: slot.description for slot_name, slot in sv.all_slots().items()}

#     return slots_dict, all_classes


def get_tier1_fields():
    tier1_cards = {}

    datacard_dict = yaml.safe_load(GENESIS_MISSION_DATA_CARD_BREAKDOWN.read_text(encoding="utf-8"))

    # creates a dict of all flattened fields whose value is if it is required or not (true/false)
    flattened_fields = flattened_tier1_fields(datacard_dict)

    # YAML dict - TODO: Can maybe skip this since not using this in tier1_GenesisCard_agent.py
    # tier1_cards["data_card_yaml"] = yaml.safe_load(GENESIS_MISSION_DATA_CARD_YAML.read_text(encoding="utf-8"))

    # MD str
    with open(GENESIS_MISSION_DATA_CARD_MD, "r") as f:
        tier1_cards["markdown_template"] = f.read()

    # Reference guide str
    with open(GENESIS_MISSION_DATA_CARD_REFERENCE, 'r') as f:
        tier1_cards["card_reference"] = f.read()

    return flattened_fields, tier1_cards, datacard_dict


def get_conditional_info(field_path: str, yaml_dict: dict) -> dict | None:
    """
    Check if a field path is inside a conditional context block.
    Returns dict with conditional info if found, None otherwise.
    
    The discriminator field is a SIBLING of the conditional alternative,
    so we need to build the path correctly.
    
    Example:
    field_path = "discoverability.authors.person.given_name"
    
    Structure:
        authors:
        type: <discriminator>
        person: <conditional alternative>
            given_name: <this field>
            
    Returns: {
        "is_conditional": True,
        "discriminator": "type",
        "discriminator_path": "discoverability.authors.type",  # Sibling of person
        "alternative": "person"
    }
    """
    parts = field_path.split(".")

    # Walk the YAML structure following the path
    current = yaml_dict
    path_to_here = []

    for i, part in enumerate(parts):
        path_to_here.append(part)

        # Navigate into the structure
        if isinstance(current, dict):
            if part in current:
                current = current[part]
            elif "value" in current and isinstance(current["value"], dict):
                # Step into value dict
                current = current["value"]
                if part in current:
                    current = current[part]
                else:
                    return None
            else:
                return None
        else:
            return None

        # Check if THIS level has conditional_context
        if isinstance(current, dict) and current.get("conditional_context") == "one_of_alternatives":
            discriminator = current.get("condition_discriminator")

            # The discriminator is a SIBLING - go back to parent level
            # parent_path is everything EXCEPT the current part
            parent_path_parts = path_to_here[:-1]  # Remove the conditional alternative from path

            # Build discriminator path: parent + discriminator name
            discriminator_path_parts = parent_path_parts + [discriminator]

            return {
                "is_conditional": True,
                "discriminator": discriminator,
                "discriminator_path": ".".join(discriminator_path_parts),
                "alternative": part,  # The current part IS the alternative name
                "parent_path": ".".join(parent_path_parts)
            }

    return None


def flattened_tier1_fields(yaml_dict: dict):
    def get_children(field):
        value = field.get("value")

        if isinstance(value, dict):
            return value
        if isinstance(value, list):
            children = {}
            for item in value:
                if isinstance(item, dict):
                    children.update(item)
            return children

        return {}

    def walk(field, path, result, in_conditional_block=False):
        children = get_children(field)

        # This is an actual field.
        if not children:
            # If inside a conditional block, mark as not required at UI level

            result[path] = field.get("required", False)
            return
            # if in_conditional_block:
            #     result[path] = False  # ← Don't validate conditionally required fields at UI level
            # else:
            #     result[path] = field.get("required", False)
            # return

        # Check if this field is a conditional alternative
        is_conditional = field.get("conditional_context") == "one_of_alternatives" 
        
        for key, child in children.items():
            child_path = f"{path}.{key}"
            if isinstance(child, dict):
                walk(child, child_path, result, is_conditional or in_conditional_block)
            else:
                result[child_path] = False

    result = {}
    for key, field_data in yaml_dict.items():
        walk(field_data, key, result)
    return result


def extract_text_from_pdf(file) -> str:
    reader = pypdf.PdfReader(file)
    chunks = []
    for page in reader.pages:
        chunks.append(page.extract_text() or "")
    return "\n".join(chunks)


def extract_text_from_docx(file) -> str:
    doc = docx.Document(file)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text_parts = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        text_parts.append(f"--- Slide {slide_number} ---")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())

    return "\n".join(text_parts)


def parse_uploaded_context_files(uploaded_files) -> str:
    combined = []

    for uploaded_file in uploaded_files:
        suffix = Path(uploaded_file["name"] if is_remote else uploaded_file.name).suffix.lower()

        try:
            if suffix == ".pdf":
                text = extract_text_from_pdf(uploaded_file["path"] if is_remote else uploaded_file)
            elif suffix == ".docx":
                text = extract_text_from_docx(uploaded_file["path"] if is_remote else uploaded_file)
            elif suffix in [".md", ".txt"]:
                with open(uploaded_file["path"] if is_remote else uploaded_file, "r", encoding="utf-8") as f:
                    text = f.read()
            elif suffix == ".pptx":
                text = extract_text_from_pptx(uploaded_file["path"] if is_remote else uploaded_file)
            else:
                continue

            if text.strip():
                combined.append(
                    f'\n\n--- Uploaded file: {uploaded_file["name"] if is_remote else uploaded_file.name} ---\n\n{text}')

        except Exception as e:
            st.warning(f'Could not parse {uploaded_file["name"] if is_remote else uploaded_file.name}: {e}')

    return "\n".join(combined).strip()


def generate_datasheet_pdf(df: pd.DataFrame, output_pdf: str):
    row = df.iloc[0]

    doc = SimpleDocTemplate(output_pdf, pagesize=letter)
    styles = getSampleStyleSheet()
    elements = []

    for section in SECTIONS:
        if section["section_idx"] != 0:
            elements.append(Paragraph(str(section["section_idx"]) + ". " + section["title"], styles["Heading1"]))
            elements.append(Spacer(1, 12))

        for q in section["questions"]:
            qid = q["id"]

            if qid not in row.index:
                continue
            if section["section_idx"] == 0 and qid not in ["project_name", "ald", "primary_contact", "data_owner", "classification"]:
                continue

            answer = row[qid]

            if pd.isna(answer) or str(answer).strip() == "":
                continue

            question_label = html.escape(str(q["label"]))
            question_style = ParagraphStyle("QuestionStyle", parent=styles["BodyText"], fontSize=12, leading=15)
            answer_text = html.escape(str(answer))

            if section["section_idx"] != 0:
                elements.append(Paragraph(f"<b>{question_label}</b>", question_style))
                elements.append(Paragraph(answer_text, styles["BodyText"]))
            else:
                if qid == "classification" and pd.notna(row["ROSY_ID"]) and pd.notna(row["ROSY_Z_NUMBER"]):
                    elements.append(Paragraph(f"<b>{question_label}</b>: {answer_text}", question_style))
                    elements.append(Spacer(1, 12))
                    elements.append(Paragraph(f"<b>ROSY ID</b>: {html.escape(str(row["ROSY_ID"]))}", question_style))
                    elements.append(Paragraph(f"<b>Z# for ROSY ID</b>: {html.escape(str(row["ROSY_Z_NUMBER"]))}", question_style))
                else:
                    elements.append(Paragraph(f"<b>{question_label}</b>: {answer_text}", question_style))

            elements.append(Spacer(1, 12))

        elements.append(Spacer(1, 18))

    doc.build(elements)


def generate_tier1_datacard(qid: int, output_file: str):
    # TODO: maybe use agent to format yaml values as per template -- but yaml is standardized so might be fine

    tier1_tbls = get_tier1_table(qid)
    flattened_fields_dict = tier1_tbls["datacard_yaml"].iloc[0].to_dict()
    
    yaml_portion = {}
    for flattened_key, value in flattened_fields_dict.items():
        keys = flattened_key.split(".")
        current = yaml_portion

        for key in keys[:-1]:
            current = current.setdefault(key, {})

        current[keys[-1]] = value

    output_path = Path(output_file).with_suffix(".md")

    yaml_content = yaml.safe_dump(
        yaml_portion,
        sort_keys=False,
        allow_unicode=True,
        default_flow_style=False,
    ).rstrip()

    # TODO: use agent to format markdown as per template
    markdown_string = str(tier1_tbls["datacard_markdown"].iloc[0, 0]).strip()
    file_content = (
        "---\n"
        f"{yaml_content}\n"
        "---\n\n"
        f"{markdown_string}\n"
    )

    output_path.write_text(file_content, encoding="utf-8")


def has_text(val: Any) -> bool:
    return isinstance(val, str) and val.strip() != ""


def is_required(q: Dict[str, Any]) -> bool:
    return q.get("required", True)


def to_snake_case(text: str) -> str:
    # Separate acronyms from normal words:
    # "HTTPServer" -> "HTTP_Server"
    text = re.sub(r"([A-Z]+)([A-Z][a-z])", r"\1_\2", text)

    # Separate lowercase letters or numbers from capitals:
    # "sensorData" -> "sensor_Data"
    text = re.sub(r"([a-z0-9])([A-Z])", r"\1_\2", text)

    # Replace whitespace, hyphens, punctuation, etc. with one underscore
    text = re.sub(r"[^A-Za-z0-9]+", "_", text)

    # Remove leading/trailing underscores and normalize case
    return text.strip("_").lower()


def section_complete(section_idx: int, row: Dict[str, Any]) -> Tuple[bool, List[str]]:
    section = SECTION_BY_IDX[section_idx]
    missing: List[str] = []

    for q in section["questions"]:
        if not is_required(q):
            continue

        qtype = q.get("type", "text")
        qid = q["id"]

        if qtype == "table":
            cols = q.get("columns", ["version", "date", "authors", "changes"])
            min_full = int(q.get("min_full_rows", 1))

            rows = parse_table_json(row.get(qid, ""))
            rows = update_table(rows, cols)
            full_rows = sum(1 for r in rows if is_full_table_row(r, cols))

            if full_rows < min_full:
                missing.append(qid)
        else:
            # skip short_project_title
            if not has_text(row.get(qid, "")) and qid != "short_project_title":
                missing.append(qid)

    # Special Section 0 rule: user must either upload context files OR fill all three fields
    if section_idx == 0:
        desc_id = "project_description"
        composition_id = "data_composition"
        urls_id = "project_url"

        has_context_file = has_text(row.get("context_files", "")) or has_text(
            row.get("context_files_text", ""))
        has_all_manual_context = all(
            has_text(row.get(qid, ""))
            for qid in [desc_id, composition_id, urls_id]
        )

        if not has_context_file and not has_all_manual_context:
            if is_remote:
                missing.append("Either input a valid file in Question 7 or complete all of Questions 8, 9, and 10.")
            else:
                missing.append("Either upload a file in Question 7 or complete all of Questions 8, 9, and 10.")

        # Check for API key configuration
        if not os.environ.get("AI_API_KEY"):
            missing.append("No API key found with this app. Set AI_API_KEY environment variable.")
        if not os.environ.get("AI_API_URL"):
            missing.append("No API base URL found with this app. Set AI_API_URL environment variable.")
        if not os.environ.get("AI_MODEL"):
            missing.append("No AI Model selected for this app. Set AI_MODEL environment variable.")

    return (len(missing) == 0, missing)


def accessed_section_idx(row: Dict[str, Any], return_type) -> int:
    """
    Section 0 must be complete before user sees other sections.
    Return first incomplete section > 0.
    - If everything complete, return 0 or last section (if max return type).
    """
    ok0, _ = section_complete(0, row)
    if not ok0:
        return 0

    for idx in ALL_SECTION_IDXS:
        if idx == 0:
            continue
        ok, _ = section_complete(idx, row)
        if not ok:
            return idx

    if return_type == "first":
        return 0
    elif return_type == "max":
        return max(ALL_SECTION_IDXS)


def parse_table_json(val: Any) -> List[Dict[str, Any]]:
    if isinstance(val, str) and val.strip():
        try:
            data = json.loads(val)
            return data if isinstance(data, list) else []
        except json.JSONDecodeError:
            return []
    return []


def update_table(rows: List[Dict[str, Any]], cols: List[str]) -> List[Dict[str, Any]]:
    out = []
    for r in rows:
        rr = {}
        for c in cols:
            rr[c] = "" if r.get(c) is None else r.get(c)
        out.append(rr)
    return out


def clear_context_file_state(qid_token: str) -> None:
    keys_to_remove = [
        f"context_files_widget__{qid_token}",
        f"context_files_uploader__{qid_token}",
        f"context_files_files_list__{qid_token}",
    ]

    for key in keys_to_remove:
        st.session_state.pop(key, None)


def apply_data_editor_diff(base_rows: List[Dict[str, Any]], diff: Dict[str, Any], cols: List[str]) -> List[Dict[str, Any]]:
    """
    base_rows: list of dict rows (original)
    diff: dict with keys edited_rows, added_rows, deleted_rows
    returns the updated base_rows (new version)
    """
    rows = update_table(list(base_rows), cols)

    edited = diff.get("edited_rows") or {}
    added = diff.get("added_rows") or []
    deleted = diff.get("deleted_rows") or []

    # deleted_rows is list of row indices --- delete from highest to lowest so indices don't shift
    for idx in sorted(deleted, reverse=True):
        if isinstance(idx, int) and 0 <= idx < len(rows):
            rows.pop(idx)

    # edited_rows is usually dict
    for idx_str, changes in edited.items():
        try:
            idx = int(idx_str)
        except Exception:
            idx = idx_str if isinstance(idx_str, int) else None
        if idx is None or not (0 <= idx < len(rows)) or not isinstance(changes, dict):
            continue
        for c in cols:
            if c in changes:
                rows[idx][c] = changes[c]

    # added_rows is usually list of dicts
    for r in added:
        if isinstance(r, dict):
            new_row = {c: r.get(c, "") for c in cols}
            rows.append(new_row)

    return update_table(rows, cols)


def is_full_table_row(r: Dict[str, Any], cols: List[str]) -> bool:
    # Checks if every column is empty --- cannot proceed with section 1
    for c in cols:
        v = r.get(c, None)
        if v is None:
            return False
        if isinstance(v, str) and v.strip() == "":
            return False
        if str(v).strip() == "":
            return False
    return True


def drop_all_empty_rows(rows: List[Dict[str, Any]], cols: List[str]) -> List[Dict[str, Any]]:
    def all_empty(r: Dict[str, Any]) -> bool:
        return all(str(r.get(c, "")).strip() == "" for c in cols)
    return [r for r in rows if isinstance(r, dict) and not all_empty(r)]


def _qid_token(active_qid: int | None) -> str:
    return "draft" if active_qid is None else f"qid{int(active_qid)}"


def apply_section_updates(section_idx: int, qid_token: str) -> Dict[str, Any]:
    updates: Dict[str, Any] = {}
    section = SECTION_BY_IDX[section_idx]
    for q in section["questions"]:
        col = q["id"]
        qtype = q.get("type", "text")
        key = f"q__{col}__{qid_token}"

        if qtype == "file_uploader":
            # Get the actual file objects to extract metadata
            files_list_key = f"{col}_files_list__{qid_token}"
            uploaded_files = st.session_state.get(files_list_key, [])

            # Store file metadata (names, types, sizes) as JSON
            if uploaded_files:
                file_metadata = [
                    {"name": f["name"], "size": f["size"]} if is_remote else {"name": f.name, "size": f.size}
                    for f in uploaded_files
                ]
                updates[col] = json.dumps(file_metadata)

                uploader_key = f"{col}_uploader__{qid_token}"
                if is_remote:
                    parsed_text = parse_uploaded_context_files(uploaded_files)
                else:
                    parsed_text = st.session_state.get(uploader_key, "")
                # Store extracted text separately
                updates["context_files_text"] = parsed_text
            else:
                continue

        elif qtype == "table":
            cols = q.get("columns", ["version", "date", "authors", "changes"])
            original_key = f"original__{col}__{qid_token}"
            widget_key = f"q__{col}__{qid_token}"

            if original_key not in st.session_state:
                continue

            base_rows = st.session_state.get(original_key, [])
            diff = st.session_state.get(widget_key, {})

            if isinstance(diff, dict) and {"edited_rows", "added_rows", "deleted_rows"} <= set(diff.keys()):
                new_rows = apply_data_editor_diff(base_rows, diff, cols)
            else:
                new_rows = diff if isinstance(diff, list) else base_rows

            new_rows = update_table(new_rows, cols)
            rows_to_save = drop_all_empty_rows(new_rows, cols)
            if not rows_to_save and base_rows:
                rows_to_save = drop_all_empty_rows(base_rows, cols)
            st.session_state[original_key] = keep_one_blank_row(new_rows, cols)
            updates[col] = json.dumps(rows_to_save)
        else:
            if key not in st.session_state:
                continue
            val = st.session_state[key]

            # check to ensure short title is not None
            if col == "short_project_title":
                if not val.strip():
                    st.error("Please fill all required fields before continuing:")
                    st.write("- Short Descriptive Title")
                    st.stop()

            # special check to see if short_project_title changed
            if col == "short_project_title" and token != "draft":
                project_id = int(qid_token[3:]) # qid3 or qid33
                # compare t1 db names, not full paths
                curr_tier1_name = get_tier1_db_path(project_id, True)
                new_tier1_name = make_tier1_db_name(val, True)

                if curr_tier1_name != new_tier1_name:
                    # make new tier2 name as well
                    tier2_db_path = get_tier2_db_path(project_id)
                    new_tier2_name = make_tier2_db_name(val, True)

                    store = get_db(get_master_db_name())

                    # if changed, check that new short title is not same as current ones
                    df = store.query(f"SELECT tier1_db_path FROM {PROJECTS_TABLE} WHERE project_id != {project_id}", True)
                    curr_proj_list = df["tier1_db_path"].str.removesuffix("_tier1.db").tolist()
                    if new_tier1_name.removesuffix("_tier1.db") in curr_proj_list:
                        st.error("Short Descriptive Title must be unique. A project with this title already exists locally.")
                        st.stop()

                    this_project = store.find(f"project_id = {project_id}", True, True)
                    this_project["tier1_db_path"] = new_tier1_name
                    this_project["tier2_db_path"] = new_tier2_name
                    store.update(this_project)
                    store.close()

                    new_tier1_path = make_tier1_db_name(val)
                    tier1_db_path = str(get_maven_dir() / curr_tier1_name)
                    shutil.copy2(tier1_db_path, new_tier1_path)
                    os.remove(tier1_db_path)

                    # check if tier2 db exists and only then copy
                    if os.path.exists(tier2_db_path):
                        new_tier2_path = make_tier2_db_name(val)
                        shutil.copy2(tier2_db_path, new_tier2_path)
                        os.remove(tier2_db_path)

                continue # don't add to updates
            if col == "project_name":
                val = val.upper()
            updates[col] = val
    return updates


def keep_one_blank_row(rows: List[Dict[str, Any]], cols: List[str]) -> List[Dict[str, Any]]:
    if rows and len(rows) > 0:
        return rows
    return [{c: "" for c in cols}]


def render_section(section_idx: int, row: Dict[str, Any], qid_token: str) -> None:
    section = SECTION_BY_IDX[section_idx]
    st.subheader(section["title"])

    if section.get("description"):
        for s in section["description"]:
            st.markdown(f"##### {s}")

    q_counter = 1
    for q in section["questions"]:
        qtype = q.get("type", "text")
        col = q["id"]
        required = is_required(q)

        q_num = f"**{q_counter}.** " if section_idx == 0 else f"**{col[1:]}** "
        q_counter += 1
        label = q_num + q["label"] + (" *" if required else " (optional)")
        caption = q.get("caption", None)

        if qtype == "file_uploader":
            if q.get("intro"):
                st.space()
                for i in q["intro"]:
                    if is_remote:
                        i = i.replace("Upload", "Input absolute path to")
                    st.markdown(f"##### {i}")

            widget_key = f"{col}_widget__{qid_token}"
            key = f"{col}_uploader__{qid_token}"

            if is_remote:
                st.write("Enter absolute path to context files -- 1 per line")
            else:
                st.write(label)

            existing_files = parse_table_json(row.get(col, ""))
            if existing_files:
                st.info("Previously uploaded file(s): " + ", ".join(f.get("name", "unknown") for f in existing_files))

            files_list_key = f"{col}_files_list__{qid_token}"
            if is_remote:
                uploaded_files = st.text_area(label, key=widget_key, height=100, label_visibility="collapsed")

                uploaded_files = [{"path": p, "name": p.name, "size": p.stat().st_size} for line in uploaded_files.splitlines() 
                                  if (line1 := line.strip()) if (p := Path(line1)).is_file()]
                st.session_state[files_list_key] = uploaded_files
            else:
                uploaded_files = st.file_uploader(label, type=["pdf", "docx", "md", "txt", "pptx"], label_visibility="collapsed",
                                                accept_multiple_files=True, key=widget_key)
                if uploaded_files:
                    # Save file objects to session state for metadata extraction in apply_section_updates
                    st.session_state[files_list_key] = uploaded_files

                combined_text = parse_uploaded_context_files(uploaded_files)
                st.session_state[key] = combined_text
                if combined_text:
                    file_str = "file" if len(uploaded_files) == 1 else "files"
                    st.success(f"Uploaded {len(uploaded_files)} {file_str}.")

        elif qtype == "table":
            cols = q.get("columns", ["version", "date", "authors", "changes"])
            original_key = f"original__{col}__{qid_token}"

            if original_key not in st.session_state:
                existing = parse_table_json(row.get(col, ""))
                existing = update_table(existing, cols)
                if not existing:
                    existing = [dict.fromkeys(cols, "")]
                    existing = update_table(existing, cols)
                existing = keep_one_blank_row(existing, cols)
                st.session_state[original_key] = existing
            else:
                st.session_state[original_key] = keep_one_blank_row(
                    update_table(st.session_state[original_key], cols), cols)
            widget_key = f"q__{col}__{qid_token}"

            st.write(label)
            if caption is not None:
                st.caption(caption)
            st.data_editor(
                st.session_state[original_key],
                num_rows="dynamic",
                width="stretch",
                key=widget_key,
                column_config={
                    "version": st.column_config.TextColumn("Version"),
                    "date": st.column_config.TextColumn("Date"),
                    "authors": st.column_config.TextColumn("Authors"),
                    "changes": st.column_config.TextColumn("Changes"),
                },
            )

        elif qtype == "button":
            key = f"q__{col}__{qid_token}"
            values = q.get("values", [])
            st.write(label)
            if caption is not None:
                st.caption(caption)
            existing_value = row.get(col, "") or None

            index = None
            if existing_value in values:
                index = values.index(existing_value)

            st.radio(label, options=values, index=index, key=key, horizontal=False, label_visibility="collapsed")

        else:
            key = f"q__{col}__{qid_token}"
            st.write(label)
            if caption is not None:
                st.caption(caption)
            q_height = q.get("height", None)
            height = q_height if q_height is not None else 130

            value = row.get(col, "") or ""
            if col == "short_project_title" and qid_token != "draft":
                project_id = int(qid_token[3:]) # qid3 or qid33
                value = get_tier1_db_path(project_id, True).removesuffix("_tier1.db")
            st.text_area(label, value=value,
                         key=key, height=height, label_visibility="collapsed")


def go_to_section(qid: int, target_idx: int) -> None:
    st.session_state.active_qid = int(qid)
    st.session_state.section_idx = int(target_idx)


def commit_section0_and_create_row(section0_row: Dict[str, Any]) -> int:
    section0_row["project_name"] = section0_row["project_name"].upper()

    short_title = section0_row["short_project_title"]
    tier1_db_path = make_tier1_db_name(short_title)
    create_tier1_db(tier1_db_path)

    tier1_db_name = make_tier1_db_name(short_title, True)
    # do not make file, but add path to the projects table
    tier2_db_name = make_tier2_db_name(short_title, True)

    updates0 = {
        q["id"]: section0_row.get(q["id"], "")
        for q in SECTION_BY_IDX[0]["questions"]
        if q["id"] != "short_project_title"
    }

    if "context_files_text" in section0_row:
        updates0["context_files_text"] = section0_row["context_files_text"]

    project_store = get_db(tier1_db_path)
    project_store.read(updates0, "Collection", table_name=DATASHEET_TABLE)
    project_store.close()

    master_store = get_db(get_master_db_name())
    df = master_store.query(f"SELECT project_id FROM {PROJECTS_TABLE} ORDER BY project_id DESC LIMIT 1;", True)
    curr_proj_id = 1 if df.empty else int(df.iloc[0, 0]) + 1

    master_dict = {"project_id": curr_proj_id, "project_name": section0_row["project_name"],
                   "tier1_db_path": tier1_db_name, "tier2_db_path": tier2_db_name, "has_moved": "no"}
    master_store.read(master_dict, "Collection", table_name=PROJECTS_TABLE)
    master_store.close()

    return curr_proj_id


def render_autofill_review(row: Dict[str, Any]) -> None:    
    meta = load_agent_meta(row)
    payload = {"fields": meta.get("fields", {})}
    summary = summarize_autofill(payload, SECTIONS)

    source = meta.get("last_source", "Heuristic Autofill")
    st.write(f'## {source.replace("_", " ").title()} Review')
    st.write(" ")

    if summary["filled"]:
        st.write("#### Autofilled fields")
        st.write("##### Note: Users can edit autofilled fields after answering clarification questions")

        for item in summary["filled"][:12]:
            st.write(f"- **`{item['qid'][1:]}`** {item['label']}")
            st.caption(item["answer"][:240] if item["answer"]
                       else "No text stored.")
    st.write(" ")
    if summary["needs_user"] or summary["unknown"]:
        st.write("#### Fields that still need user input")
        for item in (summary["needs_user"] + summary["unknown"])[:12]:
            st.write(f"- **`{item['qid'][1:]}`** {item['label']}")
    st.write(" ")
    if summary["low_confidence"]:
        st.write("#### Low-confidence autofill")
        for item in summary["low_confidence"][:8]:
            st.write(f"- **`{item['qid'][1:]}`** {item['label']}")
            if item["rationale"]:
                st.caption(item["rationale"])
    st.write(" ")
    followups = meta.get("followup_questions", [])
    if followups:
        st.write("#### Suggested follow-up prompts (Answer on next screen)")
        for question in followups:
            st.write(f"- {question}")
    else:
        st.success(
            "No follow-up prompts were generated. You can continue to the remaining sections.")


def render_followup_form(row: Dict[str, Any], qid: int) -> None:
    st.subheader("Follow-Up Clarifications")
    meta = load_agent_meta(row)
    questions = meta.get("followup_questions") or build_followup_questions(
        row,
        {"fields": meta.get("fields", {})},
        SECTIONS,
    )
    if not questions:
        st.info("No clarification prompts remain. Continue to the datasheet sections.")
        return

    st.write("##### Clarifications can be skipped if the information is unknown or not applicable.")
    for idx, question in enumerate(questions, start=1):
        key = f"followup__{qid}__{idx}"
        question = re.sub(r' \(s\d{1,2}\.\d{1,2}\)$', '', question)
        st.write(f"**{idx}.** {question.replace("Section 0", "the Project Description")}")
        st.text_area(
            label=f"{idx}. {question}",
            key=key,
            height=140,
            label_visibility="collapsed"
        )


def collect_followup_answers(qid: int, questions: List[str]) -> Dict[str, str]:
    answers: Dict[str, str] = {}
    for idx, question in enumerate(questions, start=1):
        key = f"followup__{qid}__{idx}"
        value = st.session_state.get(key, "")
        if isinstance(value, str) and value.strip():
            answers[question] = value.strip()
    return answers


def route_after_autofill(row: Dict[str, Any]) -> int:
    target = accessed_section_idx(row, "first")
    if target == 0 and section_complete(0, row)[0]:
        return 1
    return target


def run_initial_autofill_for_project(chat_agent: ChatAgent, qid: int) -> Dict[str, Any]:
    row = get_datasheet(qid)
    payload = run_initial_autofill(chat_agent, row, SECTIONS)
    updates, _ = merge_autofill_result(
        row,
        payload,
        ALL_QUESTIONS,
        previous_meta=load_agent_meta(row),
    )
    update_datasheet(qid, updates)
    return get_datasheet(qid)


def run_followup_autofill_for_project(chat_agent: ChatAgent, qid: int, clarifications: Dict[str, str]) -> Dict[str, Any]:
    row = get_datasheet(qid)

    existing = {}
    current = row.get("user_clarifications", "")
    if isinstance(current, str) and current.strip():
        try:
            existing = json.loads(current)
        except json.JSONDecodeError:
            existing = {}
    existing.update(clarifications)
    payload = run_followup_autofill(chat_agent, row, clarifications, SECTIONS)
    updates, _ = merge_autofill_result(
        row,
        payload,
        ALL_QUESTIONS,
        previous_meta=load_agent_meta(row),
    )
    updates["user_clarifications"] = json.dumps(
        existing, indent=2, sort_keys=True)
    update_datasheet(qid, updates)
    return get_datasheet(qid)


def rebuild_yaml_structure(flattened_keys_dict):
    result = {}

    for flattened_key, value in flattened_keys_dict.items():
        parts = flattened_key.split(".")
        current = result

        for part in parts[:-1]:
            current = current.setdefault(part, {})

        current[parts[-1]] = value

    return result


def ai_model_message(url: str):
    if "aiportal-api" in url:
        st.markdown(
            """
            <div style="background-color: rgba(255, 75, 75, 0.1); padding: 1rem; border: 2px solid rgb(255, 75, 75); 
            border-radius: 0.25rem; font-size: 1.15rem; font-weight: 600;">
            ⚠️ IMPORTANT: Model-Specific Data Restrictions
            
            Different AI models support different data levels. YOU are responsible for verifying that your selected model supports
            the data level of the content you send via the API.
            </div>
            """,
            unsafe_allow_html=True,
        )
        st.write(" ")
        st.markdown("""
        - <span style="font-weight:800;">AWS GovCloud Bedrock</span> and <span style="font-weight:800;">LANL IT Hosted</span> 
                    models are approved for unclassified data levels, including: 
        OPEN, CUI, LA-CP, PHI, PII, RSI, TPI, UCNI, U-NNPI, DOE 810, ECI, EAR, ITAR.
        These include the following models:
            - <span style="font-weight:800;">Claude 4.5 Sonnet</span> (AWS GovCloud Bedrock)
            - <span style="font-weight:800;">Claude Opus 4.8</span> (AWS GovCloud Bedrock)
            - <span style="font-weight:800;">ChatGPT 5.4 - Gov</span> (AWS GovCloud Bedrock)
            - <span style="font-weight:800;">Gemma 4 31B</span> (LANL IT)
            - <span style="font-weight:800;">Nemotron 3 Super 120B</span> (LANL IT)
            - <span style="font-weight:800;">ChatGPT OSS 120B</span> (LANL IT)
            - <span style="font-weight:800;">Nova Pro v1</span> (AWS GovCloud Bedrock)
        - <span style="font-weight:800;">Azure Commercial Foundry</span> models are approved for OPEN, CUI, LA-CP, PHI, PII, TPI but 
                    <span style="font-weight:800;">:red[ARE NOT approved for export controlled or nuclear information]</span>. 
                    The following models have restrictions:
            - <span style="font-weight:800;">ChatGPT 5.5</span>: :red[NO ECI, NO ITAR, NO EAR, NO UCNI, NO UNNPI, NO RSI, NO DOE 810]
            - <span style="font-weight:800;">ChatGPT 5.4 - Comm</span>: :red[NO ECI, NO ITAR, NO EAR, NO UCNI, NO UNNPI, NO RSI, NO DOE 810]
        """, unsafe_allow_html=True)
        st.write(" ")
    elif "circe-keys" in url:
        st.markdown(':red[<span style="font-weight:800;">NOTE: All models can currently handle CUI level data</span>]', unsafe_allow_html=True)


def validate_rosy(rosy_id:str, rosy_z_num:int, verify_ssl=False):
    base_url = "https://rassti.lanl.gov"
    try:
        response = requests.get(base_url, verify=verify_ssl, timeout=10)
        response.raise_for_status()
    except requests.exceptions.RequestException as e:
        st.error(f"Failed to connect to {base_url}; check network. Error: {str(e)}")
        st.stop()

    test_url = f"{base_url}/api/query/rosy?submitter_znumber={rosy_z_num}&rosy_pid={rosy_id}"
    try:
        response = requests.get(
            test_url,
            verify=verify_ssl,
            timeout=2
        )
        response.raise_for_status()
        data = response.json()

        if data.get("error"):
            st.error("Error while checking if the ROSY ID was reviewed. Review ID and Z# fields.")
            st.stop()
        if not (data.get("rosy_pid") and data.get("submitter_znumber") and data.get("review_complete")):
            st.error("Error while checking if the ROSY ID was reviewed. Review ID and Z# fields.")
            st.stop()

        return data["review_complete"]

    except requests.exceptions.ConnectionError:
        st.error(f"Connection failed: Cannot connect to {base_url}. Ensure you are on an approved network.")
        st.stop()
    except requests.exceptions.HTTPError as e:
        if e.response.status_code == 404:
            st.error(f"ROSY API not found at {base_url}. Verify this is a valid ROSY endpoint.")
            st.stop()
        else:
            st.error(f"{str(e)}. Review ID and Z# fields")
            st.stop()
    except ValueError as e:
        st.error(f"Invalid JSON response from {base_url}: {str(e)}")
        st.stop()
    except Exception:
        st.error("Error accessing ROSY; ensure you are on an approved network.")
        st.stop()


@st.dialog("Confirm Data Changes", width="medium", dismissible=False)
def confirm_unchanged_data_dialog(qid: int, local_data: str):
    st.subheader(f"The data at '{local_data}' was recently moved")

    st.write("Has the data at the local path changed since the last move?")

    yes_col, no_col = st.columns(2)
    with yes_col:
        if st.button("Unchanged", type="primary", width="stretch"):
            st.session_state.unchanged_data = True
            st.session_state.ran_change_dialog = True
            st.rerun()
    with no_col:
        if st.button("Changed", type="secondary", width="stretch"):
            st.session_state.unchanged_data = False
            st.session_state.ran_change_dialog = True
            st.rerun()


@st.dialog("Update Short Descriptive Title", width="medium", dismissible=False)
def update_short_project_title_dialog(qid: int, short_proj_name: str, hpc_campaign: str):

    st.subheader(f"The project '{short_proj_name}' already exists in '{hpc_campaign}'")
    st.write("Update the Short Descriptive Title by adding a suffix such as '_1' or entering a new short name.")

    curr_short_name = get_tier1_db_path(qid, True).removesuffix("_tier1.db")
    short_title_input = st.text_input("Short Project Title", value=curr_short_name)

    if st.button("Save", type="primary", width="stretch"):
        cleaned_name = short_title_input.strip()
        if not cleaned_name:
            st.error("Short Project Title cannot be empty.")
            st.stop()
        if curr_short_name.lower() == cleaned_name.replace(" ", "_").lower():
            st.error("Short Project Title cannot be the same.")
            st.stop()

        curr_tier1_path = get_tier1_db_path(qid)
        curr_tier2_path = get_tier2_db_path(qid)

        new_tier1_name = make_tier1_db_name(cleaned_name, True)
        new_tier2_name = make_tier2_db_name(cleaned_name, True)

        store = get_db(get_master_db_name())

        # check that new short title is not same as current local ones
        df = store.query(f"SELECT tier1_db_path FROM {PROJECTS_TABLE} WHERE project_id != {qid}", True)
        curr_proj_list = df["tier1_db_path"].str.removesuffix("_tier1.db").tolist()
        if new_tier1_name.removesuffix("_tier1.db") in curr_proj_list:
            st.error("Short Descriptive Title must be unique locally too. Please enter another title.")
            st.stop()

        this_project = store.find(f"project_id = {qid}", True, True)
        this_project["tier1_db_path"] = new_tier1_name
        this_project["tier2_db_path"] = new_tier2_name

        store.update(this_project)
        store.close()

        new_tier1_path = make_tier1_db_name(cleaned_name)
        new_tier2_path = make_tier2_db_name(cleaned_name)


        shutil.copy2(curr_tier1_path, new_tier1_path)
        os.remove(curr_tier1_path)

        shutil.copy2(curr_tier2_path, new_tier2_path)
        os.remove(curr_tier2_path)

        st.success("Short project title updated.")
        st.session_state.change_short_proj_title = False
        st.rerun()


@st.dialog("Update MAVEN Directory", width="medium")
def update_maven_dir_dialog():
    st.write("Enter a directory where MAVEN should store all metadata databases.")

    curr_dir = get_maven_dir()
    dir_input = st.text_input("MAVEN Directory", value=curr_dir, key="update_maven_dir")

    if st.button("Save", type="primary", width="stretch"):
        if save_maven_dir(dir_input):
            st.success("MAVEN Directory saved.")
            st.rerun()


@st.dialog("Update AI API Info", width="medium")
def update_api_info_dialog():
    st.subheader("Enter a new AI API Key and Base URL")

    st.write("AI API Key")
    api_key_input = st.text_input("AI API Key", value=os.environ.get("AI_API_KEY"), type="password",
                                label_visibility="collapsed", key="update_maven_api_key")

    st.write("AI Base URL")
    api_url_input = st.text_input("AI Base URL", value=os.environ.get("AI_API_URL"), label_visibility="collapsed",
                                key="update_maven_api_base_url")

    if st.button("Select New Model ➡", type="primary", width="stretch"):
        if not api_key_input.strip() or not api_url_input.strip():
                st.error("Please enter an AI API Key and Base URL")
                st.stop()

        if get_maven_dir():
            st.session_state.api_variables = [api_key_input.strip(), api_url_input.strip()]
            st.session_state.update_ai_info_screen = False
            st.session_state.select_model_screen = True
            st.rerun()


@st.dialog("Update AI API Info", width="medium", dismissible=False)
def update_ai_model_dialog():
    if st.session_state.select_model_screen:
        back_col, _ = st.columns([1, 2], width="stretch")
        with back_col:
            if st.button("⬅ Edit API Variables", width="stretch"):
                st.session_state.api_variables = []
                st.session_state.update_ai_info_screen = True
                st.session_state.select_model_screen = False
                st.rerun()

    st.subheader("Select a new AI Model to use with your API key")
    try:
        key, url = st.session_state.api_variables
        client = OpenAI(api_key=key, base_url=url, http_client=httpx.Client(verify=False))
        models = [model.id for model in client.models.list().data]
    except Exception as e:
        st.error("Error finding models with the input AI API Key and Base URL")
        st.error(e)
        st.stop()

    ai_model_message(url.lower())

    selected_model = st.selectbox("val", models, label_visibility="collapsed", index=None,
                                    key="update_ai_model_selection")

    if st.button("Save", type="primary", width="stretch"):
        if selected_model is None:
            st.error("Please select an AI Model.")
            st.stop()
        
        if get_maven_dir():
            API_KEYS_FILE.write_text(
                f"AI_API_KEY={key}\n"
                f"AI_API_URL={url}\n"
                f"AI_MODEL=openai:{selected_model}\n"
            )
            st.success("Updated AI Model")
            st.session_state.api_variables = []
            st.session_state.select_model_screen = False
            st.session_state.update_ai_info_screen = False
            load_env_keys()
            st.rerun()


@st.dialog("Context files classification", width="medium", dismissible=False)
def confirm_context_files_dialog(context_files): # add :str or :list
    st.subheader("Confirm that these context files can be sent to the selected AI Model")
    st.write(f'**Selected AI Model**: :red[{os.getenv("AI_MODEL").removeprefix("openai:")}]')
    context_files = json.loads(context_files)
    context_files = [d["name"] for d in context_files]
    files_str = "\n - ".join(context_files)
    st.write(f"**Input context file(s)**: \n - {files_str}")
    st.space()
    st.write("Note: If unsure, first submit a ROSY review to ensure it can be sent to that model")
    st.space()
    col1, col2 = st.columns(2)
    if col1.button("Yes", width="stretch"):
        st.session_state.confirm_submit_context_files = True
        st.rerun()
    if col2.button("No", width="stretch"):
        st.rerun()

# -----------------------------
# App start
# -----------------------------
st.set_page_config(page_title="MAVEN",
                   page_icon="📝", layout="wide")

st.markdown("""
<style>
[data-testid="stFileUploaderDropzone"] { min-height: 140px; align-items: center; }
# [data-testid="stFileUploaderDropzone"] button { margin-left: 12px; }

.st-key-section_btns:has(button) button { width: 100%; }
.st-key-section_btns:has(button) button > div { justify-content: flex-start !important; }
.st-key-section_btns:has(button) button p { text-align: left !important; margin: 0; width: 100%; }
</style>
""", unsafe_allow_html=True)

if st.session_state.get("_scroll_to_top"):
    st.session_state._scroll_to_top = False
    st.markdown(
        """
        <input autofocus style= "position: absolute; top: 0; left: 0; width: 0px;
        height: 0px; opacity: 0; border: 0; padding: 0; margin: 0; pointer-events: none;"/>
        """, unsafe_allow_html=True)



if "active_qid" not in st.session_state:
    st.session_state.active_qid = None
if "section_idx" not in st.session_state:
    st.session_state.section_idx = 0
if "draft_answers" not in st.session_state:
    st.session_state.draft_answers = {}
if "jumped_once" not in st.session_state:
    st.session_state.jumped_once = {}
if "draft_mode" not in st.session_state:
    st.session_state.draft_mode = False
if "screen" not in st.session_state:
    st.session_state.screen = "datasheet"
if "confirm_delete_qid" not in st.session_state:
    st.session_state.confirm_delete_qid = None
if "render_t2_extraction" not in st.session_state:
    st.session_state.render_t2_extraction = False
if "tier2_loc_dict" not in st.session_state:
    st.session_state.tier2_loc_dict = None
if "change_short_proj_title" not in st.session_state:
    st.session_state.change_short_proj_title = False
if "unchanged_data" not in st.session_state:
    st.session_state.unchanged_data = False
if "ran_change_dialog" not in st.session_state:
    st.session_state.ran_change_dialog = False
if "local_to_staging_moved" not in st.session_state:
    st.session_state.local_to_staging_moved = False
if "staging_to_campaign_moved" not in st.session_state:
    st.session_state.staging_to_campaign_moved = False
if "select_model_screen" not in st.session_state:
    st.session_state.select_model_screen = False
if "api_variables" not in st.session_state:
    st.session_state.api_variables = []
if "update_ai_info_screen" not in st.session_state:
    st.session_state.update_ai_info_screen = False
if "confirm_submit_context_files" not in st.session_state:
    st.session_state.confirm_submit_context_files = False

if get_maven_dir() is None:
    st.title("Welcome to the MAVEN App")
    st.write("Enter a space to create a directory store metadata databases and AI API variables.")

    st.write("MAVEN Directory")
    st.caption("Location where a directory will be created that stores all metadata files")
    dir_input = st.text_input("MAVEN Directory", placeholder="/path/to/maven/projects", label_visibility="collapsed")
    api_keys_exist = load_env_keys()
    if not api_keys_exist:
        st.write("AI API Key")
        st.caption("Secret token to access an LLM provider's API")
        api_key_input = st.text_input("AI API Key", placeholder="25-digit API key", label_visibility="collapsed",
                                      key="maven_api_key_new_app", type="password")
        st.write("AI Base URL")
        st.caption("Endpoint URL to access an LLM provider's API")
        api_url_input = st.text_input("AI Base URL", placeholder="https://api-key.com/", label_visibility="collapsed",
                                      key="maven_api_url_new_app")

    if st.button("Save", type="primary", width="stretch"):
        if not api_keys_exist:
            if not api_key_input.strip() or not api_url_input.strip():
                st.error("Please enter an API key and its base URL")
                st.stop()

        if save_maven_dir(dir_input):
            if not api_keys_exist:
                st.session_state.api_variables = [api_key_input.strip(), api_url_input.strip()]
                st.session_state.select_model_screen = True
                st.success("Directory and API variables saved.")
            else:
                st.success("Directory saved.")
            st.rerun()

    st.stop()

loaded_keys = load_env_keys()
if not loaded_keys:
    aim_left, aim_mid, aim_right = st.columns([0.5, 3.6, 0.5], width='stretch')
    with aim_mid:
        st.title("Welcome to the MAVEN App")
        if not st.session_state.select_model_screen:
            st.subheader("Enter AI API Key and Base URL")
            
            if not st.session_state.api_variables:
                st.write("AI API Key")
                st.caption("Secret token to access an LLM provider's API")
                if os.environ.get("AI_API_KEY"):
                    api_key_input = st.text_input("AI API Key", value=os.environ.get("AI_API_KEY"), type="password",
                                                  key="maven_api_key_existing_app", label_visibility="collapsed")
                else:
                    api_key_input = st.text_input("AI API Key", placeholder="25-digit API key", type="password",
                                                  key="maven_api_key_existing_app", label_visibility="collapsed")
                st.write("AI Base URL")
                st.caption("Endpoint (URL) to access an LLM provider's API")
                if os.environ.get("AI_API_URL"):
                    api_url_input = st.text_input("AI Base URL", value=os.environ.get("AI_API_URL"), 
                                                  key="maven_api_url_existing_app", label_visibility="collapsed")
                else:
                    api_url_input = st.text_input("AI Base URL", placeholder="https://api-key.com/", 
                                                  key="maven_api_url_existing_app", label_visibility="collapsed")
            else:
                key, url = st.session_state.api_variables
                st.write("AI API Key")
                st.caption("Secret token to access an LLM provider's API")
                api_key_input = st.text_input("AI API Key", value=key, key="maven_api_key_existing_app", 
                                              type="password", label_visibility="collapsed")
                st.write("AI Base URL")
                st.caption("Endpoint (URL) to access an LLM provider's API")
                api_url_input = st.text_input("AI Base URL", value=url, key="maven_api_url_existing_app", 
                                              label_visibility="collapsed")

            if st.button("Save", type="primary", width="stretch"):
                if not api_key_input.strip() or not api_url_input.strip():
                    st.error("Please enter an API key and its base URL")
                    st.stop()

                if get_maven_dir():
                    st.session_state.api_variables = [api_key_input.strip(), api_url_input.strip()]
                    st.session_state.select_model_screen = True
                    st.success("API variables saved.")
                    st.rerun()
        else:
            back_col, other = st.columns([1, 4], width="stretch")
            with back_col:
                if st.button("⬅ Edit API Variables", width="stretch"):
                    st.session_state.select_model_screen = False
                    st.rerun()

            try:
                key, url = st.session_state.api_variables
                client = OpenAI(api_key=key, base_url=url, http_client=httpx.Client(verify=False))
                models = [model.id for model in client.models.list().data]
            except Exception as e:
                st.error("Error finding models with the input AI API Key and Base URL")
                st.code(e)
                st.stop()
            
            st.space()
            st.subheader("Pick the AI model to use with your API key")

            ai_model_message(url.lower())

            selected_model = st.selectbox("val", models, label_visibility="collapsed", index=None,
                                        key="new_ai_model_selection")
            if st.button("Save"):
                if selected_model is None:
                    st.error("Please select an AI Model.")
                    st.stop()
                
                if get_maven_dir():
                    API_KEYS_FILE.write_text(
                        f"AI_API_KEY={key}\n"
                        f"AI_API_URL={url}\n"
                        f"AI_MODEL=openai:{selected_model}\n"
                    )
                    st.success("Saved model")
                    st.session_state.api_variables = []
                    st.session_state.select_model_screen = False
                    st.rerun()
        st.stop()

if st.session_state.select_model_screen:
    update_ai_model_dialog()

if st.session_state.update_ai_info_screen:
    update_api_info_dialog()


create_master_db()
CHAT_AGENT = configure_chat_agent()

if st.session_state.screen == "datasheet":
    # -----------------------------
    # Home screen on every load
    # -----------------------------
    if st.session_state.active_qid is None and int(st.session_state.section_idx) != 0:
        st.session_state.section_idx = 0

    if st.session_state.active_qid is None and not st.session_state.draft_mode and not st.session_state.draft_answers:

        first_l, first_m, first_r = st.columns([0.75, 5, 0.75])
        with first_m:
            st.title("MAVEN Home")
            second_l, second_r = st.columns(2)
            with second_l:
                if st.button("Edit MAVEN Directory", key="edit_maven_dir_btn", width="stretch"):
                    st.session_state.select_model_screen = False
                    st.session_state.update_ai_info_screen = False
                    update_maven_dir_dialog()
            with second_r:
                if st.button("Edit AI API Info", key="edit_ai_api_info_btn", width="stretch"):
                    st.session_state.select_model_screen = False
                    st.session_state.update_ai_info_screen = False
                    update_api_info_dialog()

            st.space()
            st.caption("Click **New Project** to start a new project or choose an existing project to continue.")

            if st.button("➕ New Project", key="new_project_top", type="primary", width="stretch"):
                st.session_state.active_qid = None
                st.session_state.section_idx = 0
                st.session_state.draft_answers = {}
                st.session_state.draft_mode = True
                st.session_state.screen = "datasheet"
                st.session_state.local_to_staging_moved = False
                st.session_state.staging_to_campaign_moved = False
                st.session_state.select_model_screen = False
                st.session_state.update_ai_info_screen = False
                st.rerun()

            rows = list_projects()
            if not rows.empty:
                st.divider()

                for r in rows.itertuples(index=True):
                    qid = int(r.project_id)
                    pname = (r.project_name or "").strip()

                    col_project, col_delete = st.columns([5, 1])

                    with col_project:
                        if st.button(pname, key=f"pick_{qid}", width="stretch"):
                            st.session_state.draft_mode = False
                            try:
                                row = get_datasheet(qid)
                                target = accessed_section_idx(row, "first")
                                go_to_section(qid, target)
                                st.session_state.jumped_once[qid] = True
                            except Exception:
                                delete_project(qid)
                                st.session_state.active_qid = None
                                st.session_state.section_idx = 0

                            if os.path.exists(get_tier2_db_path(qid)): # if t2 db exists, go to that screen
                                st.session_state.screen = "tier2"
                            elif get_tier1_table(qid, check_exists=True): # if t1 tables (not datacard) exist, go to that screen
                                st.session_state.screen = "tier1"
                            else: # else default to datasheet screens
                                st.session_state.screen = "datasheet"
                            st.session_state.local_to_staging_moved = False
                            st.session_state.staging_to_campaign_moved = False
                            st.session_state.select_model_screen = False
                            st.session_state.update_ai_info_screen = False
                            st.rerun()

                    with col_delete:
                        if st.button("🗑️", key=f"delete_{qid}", width="stretch"):
                            st.session_state.confirm_delete_qid = qid
                            st.session_state.select_model_screen = False
                            st.session_state.update_ai_info_screen = False
                            st.rerun()

                    if st.session_state.get("confirm_delete_qid") == qid:
                        st.warning(
                            f"Delete **{pname}**? This cannot be undone.")

                        confirm_col, cancel_col = st.columns(2)

                        with confirm_col:
                            if st.button("Yes, delete", key=f"confirm_delete_{qid}", type="primary", width="stretch"):
                                delete_project(qid)

                                if st.session_state.active_qid == qid:
                                    st.session_state.active_qid = None
                                    st.session_state.section_idx = 0

                                st.session_state.confirm_delete_qid = None
                                st.rerun()

                        with cancel_col:
                            if st.button("Cancel", key=f"cancel_delete_{qid}", width="stretch"):
                                st.session_state.confirm_delete_qid = None
                                st.rerun()
        st.stop()

    # -----------------------------
    # Main screen
    # -----------------------------
    active_qid = st.session_state.active_qid
    section_idx = int(st.session_state.section_idx)

    if active_qid is None:  # only Section 0 is visible
        section_idx = 0
        st.session_state.section_idx = 0
        token = _qid_token(None)

        row = dict(st.session_state.draft_answers)

        pname = (row.get("project_name") or "").strip()
        st.title(f"Metadata Collection {' for ' + pname if pname else ''}")
    else:  # Existing datasheet
        active_qid = int(active_qid)
        token = _qid_token(active_qid)

        row = get_datasheet(active_qid)
        if not row:
            st.session_state.active_qid = None
            st.session_state.section_idx = 0
            st.rerun()

        # if section 0 incomplete, force section 0
        ok0, _ = section_complete(0, row)
        if not ok0 and section_idx not in {0, REVIEW_SECTION_IDX, FOLLOWUP_SECTION_IDX}:
            section_idx = 0
            st.session_state.section_idx = 0

        if not st.session_state.jumped_once.get(active_qid, False):
            target = accessed_section_idx(row, "first")
            st.session_state.section_idx = target
            st.session_state.jumped_once[active_qid] = True

        section_idx = int(st.session_state.section_idx)

        pname = (row.get("project_name") or "").strip()
        st.title(f"Datasheet for {pname}")

    # -----------------------------
    # Progress bar
    # -----------------------------
    completed = 0
    if st.session_state.active_qid is not None:
        latest_row = get_datasheet(int(st.session_state.active_qid))
        for idx in ALL_SECTION_IDXS[1:]:
            ok, _ = section_complete(idx, latest_row)
            if ok:
                completed += 1

        st.progress(completed / (len(ALL_SECTION_IDXS) - 1))
        st.write(
            f"Progress: **{completed}/{(len(ALL_SECTION_IDXS) - 1)}** sections complete")
        st.divider()
    else:
        ok0, _ = section_complete(0, row)

    if section_idx == REVIEW_SECTION_IDX:
        render_autofill_review(row)
    elif section_idx == FOLLOWUP_SECTION_IDX:
        render_followup_form(row, int(active_qid))
    else:
        render_section(section_idx, row, token)
    st.divider()

    # -----------------------------
    # Navigation buttons
    # -----------------------------
    only_next = (
        section_idx == 0
        or st.session_state.active_qid is None
        or section_idx not in ALL_SECTION_IDXS
    )
    if only_next:
        col1, next_col, col3 = st.columns(3, width="stretch")
    else:
        back_col, save_col, next_col = st.columns(3, width="stretch")

    if not only_next:
        with back_col:
            if st.button("⬅ Back", width="stretch"):
                qid = int(st.session_state.active_qid)
                if section_idx in ALL_SECTION_IDXS:
                    updates = apply_section_updates(section_idx, token)
                    update_datasheet(qid, updates)
                    target = max(0, section_idx - 1)
                elif section_idx == FOLLOWUP_SECTION_IDX:
                    target = REVIEW_SECTION_IDX
                else:
                    target = 0
                st.session_state.section_idx = target
                st.session_state._scroll_to_top = True
                st.rerun()

    if not only_next:
        with save_col:
            if st.button("💾 Save", width="stretch"):
                qid = int(st.session_state.active_qid)
                updates = apply_section_updates(section_idx, token)
                update_datasheet(qid, updates)
                st.success("Saved changes")

    with next_col:
        is_last = (section_idx == max(ALL_SECTION_IDXS))
        if section_idx == REVIEW_SECTION_IDX:
            next_label = "Clarify ➡"
        elif section_idx == FOLLOWUP_SECTION_IDX:
            next_label = "Apply Clarifications ➡"
        elif st.session_state.active_qid is None:
            next_label = "Submit for Autofill ➡"
        else:
            next_label = "Continue to Findability Metadata ➡" if is_last else "Next ➡"

        if st.button(next_label, type="primary", width="stretch"):
            if st.session_state.active_qid is None:
                # commit section 0 only if complete
                updates0 = apply_section_updates(0, "draft")
                st.session_state.draft_answers.update(updates0)

                short_name = make_tier1_db_name(st.session_state.draft_answers.get("short_project_title", ""), True).removesuffix("_tier1.db")

                temp_store = get_db(get_master_db_name())
                df = temp_store.query(f"SELECT tier1_db_path FROM {PROJECTS_TABLE}", True)
                curr_proj_list = df["tier1_db_path"].str.removesuffix("_tier1.db").tolist()
                if short_name in curr_proj_list:
                    st.error("Short Descriptive Title must be unique. A project with this title already exists locally.")
                    st.stop()

                ok, missing = section_complete(
                    0, st.session_state.draft_answers)
                if not ok:
                    question_lookup = {q["id"]: q["label"]
                                       for q in SECTION_BY_IDX[0]["questions"]}
                    missing_labels = [
                        question_lookup.get(m, m) for m in missing]
                    st.error("Please fill all required fields before continuing:")
                    for ml in missing_labels:
                        st.write(f"- {ml}")
                    st.stop()

                if not st.session_state.confirm_submit_context_files:
                    context_files = st.session_state.draft_answers.get("context_files", "")
                    if context_files:
                        confirm_context_files_dialog(context_files)
                        st.stop()
                    else:
                        st.session_state.confirm_submit_context_files = True

                new_id = commit_section0_and_create_row(st.session_state.draft_answers)
                with st.spinner("Running URSA-assisted autofill. May take a few minutes..."):
                    row_after_autofill = run_initial_autofill_for_project(CHAT_AGENT, new_id)
                st.session_state.active_qid = new_id
                meta = load_agent_meta(row_after_autofill)
                st.session_state.section_idx = (
                    REVIEW_SECTION_IDX
                    if meta.get("fields")
                    else route_after_autofill(row_after_autofill)
                )
                st.session_state.draft_answers = {}
                st.session_state.draft_mode = False
                st.session_state.jumped_once[new_id] = True
                st.session_state._scroll_to_top = True
                st.rerun()
            elif section_idx == REVIEW_SECTION_IDX:
                meta = load_agent_meta(row)
                followups = meta.get("followup_questions", [])
                st.session_state.section_idx = (FOLLOWUP_SECTION_IDX if followups else route_after_autofill(row))
                st.session_state._scroll_to_top = True
                st.rerun()
            elif section_idx == FOLLOWUP_SECTION_IDX:
                qid = int(st.session_state.active_qid)
                questions = load_agent_meta(row).get("followup_questions", [])
                clarifications = collect_followup_answers(qid, questions)
                with st.spinner("Applying clarifications. May take a few minutes..."):
                    row_after_followup = run_followup_autofill_for_project(CHAT_AGENT, qid, clarifications)
                st.session_state.section_idx = route_after_autofill(row_after_followup)
                st.session_state._scroll_to_top = True
                st.rerun()
            else:
                qid = int(st.session_state.active_qid)
                updates = apply_section_updates(section_idx, token)
                update_datasheet(qid, updates)

                datasheet_df = get_datasheet(qid, True)
                row2 = datasheet_df.iloc[0].to_dict()
                ok, missing = section_complete(section_idx, row2)
                if not ok:
                    st.error("Please fill all required fields before continuing:")
                    if section_idx > 0:
                        for ml in missing:
                            st.write(f"- {ml[1:]}")
                    else:
                        question_lookup = {q["id"]: q["label"]
                                           for q in SECTION_BY_IDX[0]["questions"]}
                        missing_labels = [question_lookup.get(m, m) for m in missing]
                        for ml in missing_labels:
                            st.write(f"- {ml}")
                    st.stop()

                if section_idx == 0:
                    st.session_state.section_idx = 1
                    st.session_state._scroll_to_top = True
                    st.session_state.screen = "datasheet"
                    st.rerun()

                if not is_last:
                    st.session_state.section_idx = section_idx + 1
                    st.session_state._scroll_to_top = True
                    st.session_state.screen = "datasheet"
                    st.rerun()

                # go to tier 1 screen
                else:
                    full_name = to_snake_case(datasheet_df["project_name"].iloc[0].strip()) + "_datasheet.pdf"
                    generate_datasheet_pdf(datasheet_df, str(get_maven_dir() / full_name))

                    if not get_tier1_table(qid, check_exists=True): # run ai agent
                        flattened_fields, tier1_cards, datacard_dict = get_tier1_fields()
                        with st.spinner("Populating findability metadata.\nMay take a few minutes..."):
                            all_tier1_dicts = run_tier1_catalog(CHAT_AGENT, datasheet_df, datacard_dict, tier1_cards, flattened_fields)

                        tier1_db_path = get_tier1_db_path(qid)
                        store = get_db(tier1_db_path)

                        for tier1_table_name, tier1_dict in all_tier1_dicts.items():
                            store.read(tier1_dict, "Collection", tier1_table_name)
                        store.close()

                    st.session_state.screen = "tier1"
                    st.session_state.section_idx = 0
                    st.session_state._scroll_to_top = True
                    st.session_state.local_to_staging_moved = False
                    st.session_state.staging_to_campaign_moved = False
                    st.rerun()


# -----------------------------
# TIER 1 Screen
# -----------------------------
elif st.session_state.screen == "tier1":
    if st.session_state.active_qid is None:
        st.error("No project selected for displaying Findability Metadata. Choose a project on the home screen.")
        st.stop()

    qid = int(st.session_state.active_qid)
    st.title("Findability Metadata")
    st.subheader("Click the Save button at the bottom of the screen to apply any changes")

    flattened_fields, tier1_cards, datacard_dict = get_tier1_fields()

    curr_tables = get_tier1_table(qid, update=True)
    
    if not curr_tables:
        datasheet_df = get_datasheet(qid, df_return=True)
        if datasheet_df.empty: # no datasheet data so go home
            delete_project(qid)
            st.session_state.active_qid = None
            st.session_state.section_idx = 0
            st.session_state.screen = "datasheet"
            st.session_state.local_to_staging_moved = False
            st.session_state.staging_to_campaign_moved = False
            st.rerun()
        else: # run t1 agent
            with st.spinner("Populating findability metadata. May take a few minutes..."):
                all_tier1_dicts = run_tier1_catalog(CHAT_AGENT, datasheet_df, datacard_dict, tier1_cards, flattened_fields)

            tier1_db_path = get_tier1_db_path(qid)
            store = get_db(tier1_db_path)
            for tier1_table_name, tier1_dict in all_tier1_dicts.items():
                store.read(tier1_dict, "Collection", tier1_table_name)
            store.close()
            st.rerun()

    # TODO - uncomment portion that updates datacard db with new template once datacard template is stable
    # template_dict = get_breakdown_fields(datacard_dict)
    # diff_tables_dict = {}
    # new_cols = []
    # for dc_tbl, tbl_cols in template_dict.items():
    #     if dc_tbl not in curr_tables.keys():
    #         diff_tables_dict[dc_tbl] = datacard_dict[dc_tbl]
    #         new_cols.extend(tbl_cols)
    #     else:
    #         curr_tbl_cols = curr_tables[dc_tbl].columns.to_list()
    #         diff = list(set(tbl_cols) - set(curr_tbl_cols))
    #         if diff:
    #             diff_tables_dict[dc_tbl] = datacard_dict[dc_tbl]
    #             new_cols.extend(diff)

    # # if there are tier 1 md columns not in db, add them to the db and run agent to fill them
    # if diff_tables_dict:
    #     # new_tier1_fields_dict = {k:v for k,v in tier1_fields_dict.items() if k in new_cols}
    #     datasheet_df = get_datasheet(qid, True)
    #     # yaml_card_out = str(get_diana_dbs_dir()) + "/" + to_snake_case(datasheet_df["project_name"].iloc[0].strip()) + ".yaml"
    #     # yaml_card_out = get_tier1_YAML_MD_path(qid)
    #     with st.spinner("Updating tier 1 metadata catalog with new fields..."):
    #         new_tier1_dicts = run_tier1_catalog(CHAT_AGENT, datasheet_df, diff_tables_dict, tier1_cards)

    #     for tbl_name, new_tbl_data in new_tier1_dicts.items():
    #         if tbl_name in curr_tables.keys():
    #             for new_col, new_val in new_tbl_data.items():
    #                 curr_tables[tbl_name][new_col] = new_val
    #             update_tier1_table(qid, curr_tables[tbl_name])
    #         else:
    #             tier1_db = get_tier1_db_path(qid)
    #             store2 = get_db(tier1_db)
    #             store2.read(new_tbl_data, "Collection", tbl_name)
    #             curr_tables[tbl_name] = store2.get_table(tbl_name, True, True)
    #             store2.close()

    #         # yaml_card_out = get_diana_dbs_dir() / to_snake_case(datasheet_df["project_name"].iloc[0].strip()) + ".yaml"
    #         # with open(yaml_card_out, 'w') as f:
    #         #     yaml.safe_dump(all_tier1_card, stream=f, sort_keys=False)

    with st.form(f"tier1_metadata_form_{qid}"):
        updated_values = {}
        actual_field_reqs = {}
        if "datacard_yaml" in curr_tables.keys():
            tbl_name = "datacard_yaml"
            tbl_df = curr_tables[tbl_name]
            updated_values[tbl_name] = {}

            tbl_data = tbl_df.iloc[0].to_dict()

            # dont show the dsi metadata column on the form
            dsi_key = next(iter(tbl_data))
            updated_values[tbl_name][dsi_key] = tbl_data.pop(dsi_key)

            actual_yaml_struct = rebuild_yaml_structure(tbl_data)

            def render_yaml_portion(data, path=(), depth=0):

                for key, value in data.items():
                    current_path = (*path, key)
                    widget_key = ".".join(current_path)

                    # Create a blank column to indent the entire row.
                    if depth == 0:
                        content = st.container()
                    else:
                        _, content = st.columns([depth, 12])

                    if isinstance(value, dict):
                        with content:
                            if depth == 0:
                                st.write(f"### {key}")
                                st.caption(datacard_dict[key]["description"])
                            st.write(f"**{key}**")

                        render_yaml_portion(
                            value,
                            path=current_path,
                            depth=depth + 1,
                        )
                    else:
                        with content:
                            if depth > 0: # actual child fields
                                supports_key = "supports_" + current_path[0]
                                field_req = ""

                                # Check if field is part of a conditional block
                                conditional_info = get_conditional_info(widget_key, datacard_dict)  

                                if conditional_info and conditional_info["is_conditional"]:
                                    # Get discriminator value to see if this alternative is selected
                                    discriminator_path = conditional_info["discriminator_path"]
                                    discriminator_value = tbl_data.get(discriminator_path, "")
                                    
                                    if discriminator_value == conditional_info["alternative"]:
                                        # This alternative IS selected - show as required if needed
                                        if tbl_data[supports_key].lower() == "yes":
                                            # Check the raw schema for the actual required value
                                            # (since flattened_fields marks it as False)
                                            if flattened_fields[widget_key]:
                                                field_req = " * (selected)"
                                                actual_field_reqs[widget_key] = True
                                            else:
                                                field_req = " (selected)"
                                                actual_field_reqs[widget_key] = False
                                    else:
                                        # This alternative is NOT selected
                                        field_req = " (not selected)"
                                        actual_field_reqs[widget_key] = False

                                elif tbl_data[supports_key].lower() == "yes" and flattened_fields[widget_key]:
                                    field_req = " *"
                                    actual_field_reqs[widget_key] = True
                                else:
                                    actual_field_reqs[widget_key] = False

                                st.write(key + field_req)

                                # TODO: decide whether to include description for each field too
                                field_dict = datacard_dict
                                for part in current_path[:-1]:
                                    field_dict = field_dict[part]["value"]
                                field_dict = field_dict[key]
                                st.caption(field_dict["description"])
                                if "type" not in field_dict:
                                    updated_values[tbl_name][widget_key] = st.text_input(
                                        label=key,
                                        value="" if pd.isna(value) else str(value),
                                        key="field:" + widget_key,
                                        label_visibility="collapsed"
                                    )
                                elif field_dict["type"] == "radio":
                                    options = field_dict["options"]
                                    default_index = options.index(str(value).capitalize()) if str(value) in options else None
                                    updated_values[tbl_name][widget_key] = st.radio(
                                        "radio_label",
                                        options,
                                        index=default_index,
                                        key="field:" + widget_key,
                                        label_visibility="collapsed"
                                    )
                                elif field_dict["type"] == "dropdown":
                                    options = list(field_dict["options"])
                                    # saved_permission = field_dict["access_permissions"].iloc[0] if not locations_tbl.empty else None
                                    updated_values[tbl_name][widget_key] = st.selectbox(
                                        "Enter",
                                        options=options,
                                        index=options.index(value) if value is not None and value in options else None,
                                        key="field:" + widget_key,
                                        label_visibility="collapsed"
                                    )
                                else:
                                    st.error(f"Unsupported type: {field_dict['type']}")
                                    st.stop()
                            else: # display "support_" keys differently
                                actual_field_reqs[widget_key] = True
                                st.write(f"### {key} *")
                                st.caption(datacard_dict[key]["description"])

                                options = ["Yes", "No"]
                                default_index = options.index(str(value).capitalize()) if str(value) in options else None
                                updated_values[tbl_name][widget_key] = st.radio(
                                    "radio label",
                                    options,
                                    index=default_index,
                                    key="field:" + widget_key,
                                    label_visibility="collapsed", 
                                    )
            render_yaml_portion(actual_yaml_struct)
        
        if "datacard_markdown" in curr_tables.keys():
            tbl_name = "datacard_markdown"
            tbl_df = curr_tables[tbl_name]
            updated_values[tbl_name] = {}
            
            tbl_data = tbl_df.iloc[0].to_dict()

            # dont show the dsi metadata column on the form
            dsi_key = next(iter(tbl_data))
            updated_values[tbl_name][dsi_key] = tbl_data.pop(dsi_key)

            markdown_col_key = next(iter(tbl_data))

            st.write("### Markdown portion")
            st.caption("Carefully review model-generated text in this text block")
            actual_field_reqs[markdown_col_key] = True
            updated_values[tbl_name][markdown_col_key] = st.text_area(
                "enter",
                height = 750,
                value=tbl_data[markdown_col_key],
                key="datacard_markdown_portion",
                label_visibility="collapsed"
            )

        save_col, next_col = st.columns(2)
        with save_col:
            submitted = st.form_submit_button("Save Metadata", key=f"save_tier1_{qid}", width="stretch")
        with next_col:
            next_clicked = st.form_submit_button(
                "Continue to AI-Ready Metadata ➡", key=f"next_tier1_{qid}", width="stretch", type="primary")

        if submitted or next_clicked:
            total_error = ""
            if "datacard_yaml" in updated_values.keys():
                missing_support = []
                missing_fields = {}
                for col, val in updated_values["datacard_yaml"].items():
                    top_key = col.split(".", 1)[0]
                    if f"supports_{top_key}" in actual_field_reqs.keys():
                        req_field_missing = (updated_values["datacard_yaml"][f"supports_{top_key}"].lower() == "yes" and 
                                          actual_field_reqs[col] and not str(val).strip()
                                        )
                        if req_field_missing:
                            if top_key in missing_fields.keys():
                                missing_fields[top_key].append(col.replace(".", " -> "))
                            else:
                                missing_fields[top_key] = [col.replace(".", " -> ")]
                    elif not str(val).strip(): # check if missing support fields
                        missing_support.append(col)

                if missing_support:
                    total_error += "\n- ".join(missing_support) + "\n\n"
                if missing_fields:
                    for top_key, field_list in missing_fields.items():
                        total_error += f"**{top_key}**:\n- " + "\n- ".join(field_list) + "\n\n"
            if "datacard_markdown" in updated_values.keys():
                markdown_col_key = next(iter(tbl_data))
                markdown_field_value = next(iter(updated_values["datacard_markdown"].values()), None)
                if markdown_field_value is None or not str(markdown_field_value).strip():
                    total_error += "\n\n" + "- Markdown text area"

            if total_error != "":
                st.error(f"Please complete all these required metadata fields before continuing:\n\n{total_error}")
                st.stop()
            else:
                for tbl_name, new_values in updated_values.items():
                    try:
                        tier1_df = pd.DataFrame([new_values])
                        update_tier1_table(qid, tier1_df)
                    except Exception:
                        tier1_db = get_tier1_db_path(qid)
                        store = get_db(tier1_db)
                        store.read(new_values, "Collection", tbl_name)
                        store.close()

                if next_clicked:
                    store = get_db(get_master_db_name())
                    short_proj_name = store.query(f"SELECT tier1_db_path FROM {PROJECTS_TABLE} WHERE project_id = {qid}", 
                                            True).iloc[0,0].removesuffix("_tier1.db")

                    datacard_name = short_proj_name + "_genesis_datacard_v1.2.md"
                    generate_tier1_datacard(qid, str(get_maven_dir() / datacard_name))

                    st.session_state.screen = "tier2"
                    st.session_state.section_idx = 0
                    st.session_state._scroll_to_top = True
                    st.session_state.render_t2_extraction = False
                    st.session_state.local_to_staging_moved = False
                    st.session_state.staging_to_campaign_moved = False
                    st.rerun()
                else:
                    st.success("Findability Metadata Updated.")


# -----------------------------
# TIER 2 + DSI MOVE Screen
# -----------------------------
elif st.session_state.screen == "tier2":
    if st.session_state.active_qid is None:
        st.error("No project selected for extracting AI-Ready Metadata. Please choose a project on the home screen.")
        st.stop()

    qid = int(st.session_state.active_qid)
    tier2_db_path = get_tier2_db_path(qid)
    create_tier2_db(tier2_db_path)

    st.title("AI-Ready Metadata")

    show_tier2_extraction = (st.session_state.render_t2_extraction and st.session_state.tier2_loc_dict is not None)

    if not show_tier2_extraction:
        st.subheader("Enter data locations and HPC information to enable data extraction and movement")

        tier2_store = get_db(tier2_db_path)
        locations_tbl = tier2_store.get_table(TIER_2_TABLE, True)
        tier2_store.close()

        updated_tier2_dict = {}
        if is_remote:
            updated_tier2_dict["local_data_path"] = "N/A"
            updated_tier2_dict["username"] = "N/A"
            updated_tier2_dict["hpc_system"] = "N/A"
        else:
            st.write("Local Data Location")
            st.caption("Absolute path to data on current filesystem. Ex: /users/my_user/my_proj_data/")
            updated_tier2_dict["local_data_path"] = st.text_input("Enter", key=f"local_data_path_{qid}",
                                    value=locations_tbl["local_data_path"].iloc[0] if not locations_tbl.empty else "",
                                    label_visibility="collapsed")

            l1, r1 = st.columns(2)
            with l1:
                st.write("Username")
                st.caption("Username to access the HPC System. Ex: ssh **username**@hpc_system:/path/")
                updated_tier2_dict["username"] = st.text_input("Enter", key=f"hpc_username_{qid}",
                                    value=locations_tbl["username"].iloc[0] if not locations_tbl.empty else "",
                                    label_visibility="collapsed")
            with r1:
                st.write("HPC System")
                st.caption("HPC system to access. **Specify the transfer node, not head node**. Ex: ssh username@**hpc_system**:/path/")
                updated_tier2_dict["hpc_system"] = st.text_input("Enter", key=f"hpc_system_{qid}",
                                    value=locations_tbl["hpc_system"].iloc[0] if not locations_tbl.empty else "",
                                    label_visibility="collapsed")

        l2, r2 = st.columns(2)
        with l2:
            if is_remote:
                st.write("Staging Location")
                st.caption("Absolute path to directory where data is currently staged")
            else:
                st.write("HPC Staging Location")
                st.caption("Absolute path to HPC directory where data and metadata will be temporarily staged")
            updated_tier2_dict["hpc_staging_space"] = st.text_input("Enter", key=f"hpc_staging_{qid}",
                                value=locations_tbl["hpc_staging_space"].iloc[0] if not locations_tbl.empty else "",
                                label_visibility="collapsed")
        with r2:
            if is_remote:
                st.write("Campaign Location")
                st.caption("Absolute path to directory where data and metadata will be permanently stored")
            else:
                st.write("HPC Campaign Location")
                st.caption("Absolute path to HPC directory where data and metadata will be permanently stored")
            updated_tier2_dict["hpc_campaign_space"] = st.text_input("Enter", key=f"hpc_campaign_{qid}",
                                value=locations_tbl["hpc_campaign_space"].iloc[0] if not locations_tbl.empty else "",
                                label_visibility="collapsed")

        l3, r3 = st.columns(2)
        with l3:
            st.write("User Group (Optional)")
            st.caption("User group to share campaign data folder with. Ex: my_user_group")
            updated_tier2_dict["user_group"] = st.text_input("Enter", key=f"user_group_{qid}",
                                value=locations_tbl["user_group"].iloc[0] if not locations_tbl.empty 
                                and not pd.isna(locations_tbl["user_group"].iloc[0]) else "",
                                label_visibility="collapsed")
        with r3:
            st.write("Data Access Permissions (Optional)")
            st.caption("Set access permissions for all files in the campaign data folder")

            PERMISSION_OPTIONS = {
                "750": "Owner full access; group read/execute; others no access",
                "700": "Owner full access; everyone else no access",
                "755": "Owner full access; everyone else read/execute",
                "770": "Owner and group full access; others no access",
                "775": "Owner/group full access; others read/execute",
                "644": "Owner read/write; everyone else read-only",
                "640": "Owner read/write; group read-only; others no access",
                "600": "Owner read/write; everyone else no access",
                "664": "Owner/group read/write; others read-only",
                "660": "Owner and group read/write; others no access",
                "444": "Everyone read-only",
                "400": "Owner read-only; everyone else no access",
                "666": "Everyone read/write — generally unsafe",
                "777": "Everyone full access — generally unsafe",
            }
            saved_permission = locations_tbl["access_permissions"].iloc[0] if not locations_tbl.empty else None
            updated_tier2_dict["access_permissions"] = st.selectbox(
                "Enter", key=f"access_permissions_{qid}",
                options=PERMISSION_OPTIONS,
                index=list(PERMISSION_OPTIONS).index(saved_permission) if saved_permission is not None else None,
                format_func=lambda mode: (f"{mode}: {PERMISSION_OPTIONS[mode]}"),
                label_visibility="collapsed"
            )

        l4, r4 = st.columns(2)
        with l4:
            st.write("DIANA Catalog Endpoint")
            st.caption("Absolute path to HPC directory where this project will be registered in the DIANA catalog")
            updated_tier2_dict["diana_endpoint"] = st.text_input("Enter", key=f"diana_endpoint_{qid}",
                                value=locations_tbl["diana_endpoint"].iloc[0] if not locations_tbl.empty else "",
                                label_visibility="collapsed")
        with r4:
            st.write("Contact Email Address")
            st.caption("Contact email for questions about this project in the DIANA catalog.")
            updated_tier2_dict["contact_email"] = st.text_input("Enter", key=f"contact_email_{qid}",
                                value=locations_tbl["contact_email"].iloc[0] if not locations_tbl.empty else "",
                                label_visibility="collapsed")

        st.write(f"Optionally submit the generated datasheet in `{get_maven_dir()}` for a ROSY review and register it here.")
        l5, r5 = st.columns(2)
        temp_df = get_datasheet(qid)
        with l5:
            st.write("ROSY ID (Optional)")
            st.caption("Enter a valid ROSY ID")
            rosy_id_input = st.text_input("Enter", key=f"rosy_id_{qid}", 
                                value=temp_df["ROSY_ID"] if not pd.isna(temp_df["ROSY_ID"]) else "",
                                label_visibility="collapsed") 
        with r5:
            st.write("Z# (Optional)")
            st.caption("Enter the Z-Number of the person who submitted to ROSY")
            rosy_z_num_input = st.number_input("Enter", key=f"rosy_z_num_{qid}", min_value=100000, step=1, max_value=999999,
                                value=int(temp_df["ROSY_Z_NUMBER"]) if not pd.isna(temp_df["ROSY_Z_NUMBER"]) else None,
                                label_visibility="collapsed")


        if st.button("Save & Extract Metadata ➡", key=f"save_tier2_{qid}", width="stretch", type="primary"):
            missing_fields = [col for col, value in updated_tier2_dict.items() 
                              if not str(value).strip() and col not in ["user_group", "access_permissions"]]
            if missing_fields:
                required_names_dict = {
                    "local_data_path": "Local Data Location",
                    "username": "Username",
                    "hpc_system": "HPC System",
                    "hpc_staging_space": "HPC Staging Location",
                    "hpc_campaign_space": "HPC Campaign Location",
                    "diana_endpoint": "DIANA Catalog Endpoint",
                    "contact_email": "Contact Email Address"
                }
                st.error("Please fill all fields before continuing:\n- " +
                         "\n- ".join(required_names_dict.get(m, m) for m in missing_fields))
                st.stop()

            local_data_input = updated_tier2_dict["local_data_path"].strip()
            local_path = Path(local_data_input)

            hpc_staging_input = updated_tier2_dict["hpc_staging_space"].strip()
            if "scratch" not in hpc_staging_input.lower():
                st.error("HPC Staging Location must be in the 'scratch' cluster.")
                st.stop()
            hpc_staging_path = Path(hpc_staging_input)

            hpc_campaign_input = updated_tier2_dict["hpc_campaign_space"].strip()
            if "campaign" not in hpc_campaign_input.lower():
                st.error("HPC Campaign Location must be in the 'campaign' cluster.")
                st.stop()
            hpc_campaign_path = Path(hpc_campaign_input)

            diana_endpoint_input = updated_tier2_dict["diana_endpoint"].strip()
            if "campaign" not in diana_endpoint_input.lower():
                st.error("DIANA Endpoint must be in the 'campaign' cluster.")
                st.stop()
            diana_endpoint_path = Path(diana_endpoint_input)

            username_input = updated_tier2_dict["username"].strip()
            hpc_system_input = updated_tier2_dict["hpc_system"].strip()

            user_group_input = str(updated_tier2_dict["user_group"]).strip()

            t1_store = get_db(get_tier1_db_path(qid))
            # only add rosy id and z num to datasheet table if both complete
            valid_rosy_id = rosy_id_input is not None and rosy_id_input.strip()
            if valid_rosy_id and rosy_z_num_input is not None:
                reviewed_rosy_id = validate_rosy(rosy_id_input, int(rosy_z_num_input))
                if reviewed_rosy_id:
                    t1_store.query(f"UPDATE {DATASHEET_TABLE} SET ROSY_ID = ?, ROSY_Z_NUMBER = ?", params=(rosy_id_input, rosy_z_num_input))
                else:
                    st.error("This ROSY ID has not been reviewed yet and cannot be associated with this datasheet. " \
                            "Please try again later or clear the ROSY fields for now.")
                    st.stop()
            elif (valid_rosy_id and rosy_z_num_input is None) or (not valid_rosy_id and rosy_z_num_input is not None):
                st.error("If registering ROSY review, enter both ID and associated Z#. Cannot only save one and not the other.")
                st.stop()
            t1_store.close()

            # check if t2 table already has data and if it's same as current inputs
            if not locations_tbl.empty:
                existing_loc_dict = locations_tbl.iloc[0].to_dict()
                if all(str(existing_loc_dict[k]).strip() == str(updated_tier2_dict[k]).strip() for k in updated_tier2_dict.keys()):
                    st.session_state.render_t2_extraction = True
                    st.session_state.tier2_loc_dict = existing_loc_dict
                    st.rerun()

            if local_data_input.lower() == "n/a":
                if not (hpc_staging_path.is_absolute() and hpc_staging_path.is_dir()):
                    st.error("HPC Staging Location must be an absolute path to the data you can access")
                    st.stop()
                try:
                    next(hpc_staging_path.iterdir(), None)
                    with tempfile.NamedTemporaryFile(dir=hpc_staging_path, delete=True):
                        pass
                except (PermissionError, OSError):
                    st.error("HPC Staging Location must be a directory you can access")
                    st.stop()

                if not (hpc_campaign_path.is_absolute() and hpc_campaign_path.is_dir()):
                    st.error("HPC Campaign Location must be an absolute path to a directory you can access")
                    st.stop()
                try:
                    next(hpc_campaign_path.iterdir(), None)
                except (PermissionError, OSError):
                    st.error("HPC Campaign Location must be a directory you can access")
                    st.stop()

                if not (diana_endpoint_path.is_absolute() and diana_endpoint_path.is_dir()):
                    st.error("DIANA Endpoint must be an absolute path to a directory you can access")
                    st.stop()
                try:
                    next(diana_endpoint_path.iterdir(), None)
                except (PermissionError, OSError):
                    st.error("DIANA Endpoint must be a directory you can access")
                    st.stop()
                # set username
                updated_tier2_dict["username"] = getpass.getuser()
                updated_tier2_dict["hpc_system"] = socket.getfqdn()

                import grp
                try:
                    grp.getgrnam(user_group_input)
                except KeyError:
                    st.error("User Group does not exist on this HPC")
                    st.stop()
            else:
                if not (local_path.is_absolute() and local_path.is_dir() and os.access(local_path, os.R_OK | os.W_OK | os.X_OK)):
                    st.error("Local Data Location must be an absolute path to data you can access")
                    st.stop()
                if not hpc_staging_path.is_absolute():
                    st.error("HPC Staging Location must be an absolute path to a directory on HPC")
                    st.stop()
                if not hpc_campaign_path.is_absolute():
                    st.error("HPC Campaign Location must be an absolute path to a directory on HPC")
                    st.stop()
                if not diana_endpoint_path.is_absolute():
                    st.error("DIANA Endpoint must be an absolute path to a directory on HPC")
                    st.stop()

                num_prompts = 5 if user_group_input else 4
                print(f" \nPassword prompt 1/{num_prompts}: validating HPC access")
                with st.spinner(f"Validating HPC field inputs — check the terminal for 1/{num_prompts} password prompts..."):
                    cmd = ["ssh", f"{username_input}@{hpc_system_input}", "echo", "test"]
                    result = subprocess.run(cmd, capture_output=True, text=True, check=False)
                if result.returncode != 0:
                    print("View error on app")
                    st.error("Invalid HPC username or system name")
                    st.code(result.stderr)
                    st.stop()

                print(f" \nPassword prompt 2/{num_prompts}: validating HPC staging location")
                with st.spinner(f"Validating HPC field inputs — check the terminal for 2/{num_prompts} password prompts..."):
                    cmd = ["ssh", f"{username_input}@{hpc_system_input}", f'cd "{hpc_staging_input}" && pwd && ls']
                    result = subprocess.run(cmd, capture_output=True, text=True, check=False)
                if result.returncode != 0:
                    print("View error on app")
                    st.error("HPC Staging Location does not exist or is not accessible for this user")
                    st.code(result.stderr)
                    st.stop()

                print(f" \nPassword prompt 3/{num_prompts}: validating HPC campaign location")
                with st.spinner(f"Validating HPC field inputs — check the terminal for 3/{num_prompts} password prompts..."):
                    cmd = ["ssh", f"{username_input}@{hpc_system_input}", f'cd "{hpc_campaign_input}" && pwd && ls']
                    result = subprocess.run(cmd, capture_output=True, text=True, check=False)
                if result.returncode != 0:
                    print("View error on app")
                    st.error("HPC Campaign Location does not exist or is not accessible for this user")
                    st.code(result.stderr)
                    st.stop()

                print(f" \nPassword prompt 4/{num_prompts}: validating DIANA endpoint")
                with st.spinner(f"Validating HPC field inputs — check the terminal for 3/{num_prompts} password prompts..."):
                    cmd = ["ssh", f"{username_input}@{hpc_system_input}", f'cd "{diana_endpoint_input}" && pwd && ls']
                    result = subprocess.run(cmd, capture_output=True, text=True, check=False)
                if result.returncode != 0:
                    print("View error on app")
                    st.error("DIANA Endpoint does not exist or is not accessible for this user")
                    st.code(result.stderr)
                    st.stop()

                if user_group_input:
                    print(f" \nPassword prompt 5/{num_prompts}: validating user group is valid")
                    with st.spinner("Validating user group input — check the terminal for 4/{num_prompts} password prompts..."):
                        cmd = ["ssh", f"{username_input}@{hpc_system_input}", shlex.join(["getent", "group", user_group_input])]
                        result = subprocess.run(cmd, capture_output=True, text=True, check=False)
                    if result.returncode != 0:
                        print("View error on app")
                        st.error("User Group does not exist on this HPC")
                        st.code(result.stderr)
                        st.stop()
                print(" \nGo back to app")

            col_names = locations_tbl.columns.tolist()
            new_values = list(updated_tier2_dict.values())
            store = get_db(tier2_db_path)
            if locations_tbl.empty:
                new_tier2_dict = row_dict = dict(zip(col_names, new_values))
                store.read(new_tier2_dict, "Collection", table_name=TIER_2_TABLE)
            else:
                query = (f'UPDATE {TIER_2_TABLE} SET ' + ", ".join(f'"{k}" = ?' for k in col_names))
                store.query(query, params=new_values)
            store.close()

            st.session_state.render_t2_extraction = True
            st.session_state.tier2_loc_dict = updated_tier2_dict
            st.session_state.local_to_staging_moved = False
            st.session_state.staging_to_campaign_moved = False
            st.rerun()




    # DSI MOVE PORTION
    else:
        st.divider()
        locations_dict: dict[str, str] = st.session_state.tier2_loc_dict
        local_data = locations_dict["local_data_path"]
        username = locations_dict["username"]
        hpc_name = locations_dict["hpc_system"]
        hpc_staging = locations_dict["hpc_staging_space"]
        hpc_campaign = locations_dict["hpc_campaign_space"]
        user_group = locations_dict["user_group"]
        access_permissions_code = locations_dict["access_permissions"]
        diana_endpoint = locations_dict["diana_endpoint"]
        contact_email = locations_dict["contact_email"]

        back_col, other = st.columns([1, 4], width="stretch")
        with back_col:
            if st.button("⬅ Edit HPC Info", width="stretch"):
                st.session_state.render_t2_extraction = False
                st.session_state.tier2_loc_dict = None
                st.session_state.ran_change_dialog = False
                st.session_state.local_to_staging_moved = False
                st.session_state.staging_to_campaign_moved = False
                st.rerun()

        st.subheader(f"Extracting AI-Ready metadata from: **{local_data if local_data.lower() != 'n/a' else hpc_staging }**")
        st.caption("Run scripts here to create index of files (dircrawl), extract data types, and file-level metadata unique to each dataset.")

        st.space()

        if local_data.lower() == "n/a":
            st.subheader("Click the button to move data from staging space -> campaign space.")
        else:
            st.subheader("Click the button to move data from local space -> HPC staging -> HPC campaign")
        st.space()
        dsi_move_btn = st.button("DSI Move", key=f"dsi_move_{qid}", width="stretch", type="primary")

        if dsi_move_btn:
            temp_master_store = get_db(get_master_db_name())
            has_moved = temp_master_store.query(f"SELECT has_moved FROM {PROJECTS_TABLE} WHERE project_id = {qid}", True).iloc[0, 0]
            temp_master_store.close()

            if has_moved.lower() != "no" and not st.session_state.ran_change_dialog:
                confirm_unchanged_data_dialog(qid, has_moved)
                st.stop()


            # if local data has been moved in this session and user hasn't changed screens, then check for existing folders already happened earlier
            if not st.session_state.local_to_staging_moved:
                # if short_project_title is in campaign dir, force user to rename
                short_proj_name = get_tier1_db_path(qid, True)
                short_proj_name = short_proj_name.removesuffix("_tier1.db")
                if st.session_state.change_short_proj_title:
                    update_short_project_title_dialog(qid, short_proj_name, hpc_campaign)
                    st.stop()
                else:
                    if local_data.lower() == "n/a":
                        # find existing campaign dirs using if already on hpc
                        hpc_path = Path(hpc_campaign)
                        try:
                            curr_folders = {p.name: p.owner() for p in hpc_path.iterdir() if p.is_dir()}
                        except Exception:
                            st.error("HPC staging path does not exist. Go to previous screen and enter a valid path.")
                            st.stop()

                    else:
                        print(" \nPassword prompt 1/1: checking existing folders on campaign")
                        with st.spinner("Checking HPC Campaign for projects — check the terminal for 1 password prompt..."):
                            cmd = ["ssh", f"{username}@{hpc_name}",
                                f'find {shlex.quote(hpc_campaign)} -mindepth 1 -maxdepth 1 -type d -exec stat -c "%U %n" {{}} \\;']
                            result = subprocess.run(cmd, capture_output=True, text=True, check=False)
                        print(" \nGo back to app")
                        if result.returncode != 0:
                            st.error("Error checking existing folders in HPC Campaign Location")
                            st.code(result.stderr)
                            st.stop()

                        curr_folders = {PurePosixPath(path).name: owner for line in result.stdout.strip().splitlines() if line.strip()
                                        for owner, path in [line.split(" ", 1)]}

                    # check if short_proj_name is in campaign folder
                    if short_proj_name in curr_folders:
                        # if current username matches the username of person owning data on remote, notify them it will overwrite
                        #TODO?: maybe allow the user to not overwrite with a stop button
                        retrieved_username = curr_folders[short_proj_name]
                        if username == retrieved_username:
                            st.warning("A project with the same name exists on campaign, but it is owned by this user. This move will overwrite that data")
                        else:
                            st.session_state.change_short_proj_title = True
                            update_short_project_title_dialog(qid, short_proj_name, hpc_campaign)
                            st.stop()

            # change current current working dir to the maven_dir so all dbs can be loaded normally
            os.chdir(str(get_maven_dir()))

            # NOTE: for now running index on all data, but eventually will use filecrawled table in t2.db as the index

            # create duplicate of current t2 db without "_tier2" so only short_project.db is name -- move this db, not actual t2
            t2_db_name = get_tier2_db_path(qid, True)
            t1_db_name = get_tier1_db_path(qid, True)
            proj_name = t2_db_name.removesuffix("_tier2.db")
            temp_t2_db_name = proj_name + ".db"
            shutil.copy2(t2_db_name, temp_t2_db_name)

            # add col in tier 1 md table that is absolute path to tier 2 db on campaign. Use the temp name (actually being moved to campaign)
            t2_db_campaign_path = os.path.join(hpc_campaign, proj_name, temp_t2_db_name)
            t1_store = get_db(t1_db_name)
            datasheet_df = t1_store.get_table(DATASHEET_TABLE, True)

            # create datasheet pdf in local data loc (or hpc staging if already starting from there)
            full_name = datasheet_df["project_name"].iloc[0].strip() + " DATASHEET.pdf"
            if local_data.lower() == "n/a":
                # hpc_staging is where the datasheet should be stored
                new_datasheet_loc = os.path.join(hpc_staging, full_name)
                new_tier1_dc_loc = os.path.join(hpc_staging, proj_name + "_genesis_datacard_v1.2.md")
                data_folder_name = Path(hpc_staging).name
            else:
                new_datasheet_loc = os.path.join(local_data, full_name)
                new_tier1_dc_loc = os.path.join(local_data, proj_name + "_genesis_datacard_v1.2.md")
                data_folder_name = Path(local_data).name
            generate_datasheet_pdf(datasheet_df, new_datasheet_loc)
            generate_tier1_datacard(qid, new_tier1_dc_loc)

            datasheet_campaign_path = os.path.join(hpc_campaign, proj_name, data_folder_name, full_name)
            datacard_campaign_path = os.path.join(hpc_campaign, proj_name, data_folder_name, proj_name + "_genesis_datacard_v1.2.md")
            if FILE_POINTERS_TABLE in t1_store.list(True):
                t1_store.query(f"UPDATE {FILE_POINTERS_TABLE} SET datsheet_file = ?, datacard_file = ?, tier2_db_path = ?;", 
                               params=(datasheet_campaign_path, datacard_campaign_path, t2_db_campaign_path))
            else:
                file_pointers_dict = {"datsheet_file": datasheet_campaign_path, 
                                      "datacard_file": datacard_campaign_path, 
                                      "tier2_db_path": t2_db_campaign_path}
                t1_store.read(file_pointers_dict, "Collection", FILE_POINTERS_TABLE)
            t1_store.close()

            skip_index = st.session_state.unchanged_data

            # diana federation endpoint entry for this project
            fed_line = f"HPC,{hpc_name},{os.path.join(hpc_campaign, t1_db_name)},data,{username},{contact_email},{datetime.now(UTC).time()}"

            if local_data.lower() == "n/a":
                result = subprocess.run(["module avail conduit"], shell=True, executable="/bin/bash", capture_output=True)
                if "conduit/conduit-x86_64" in str(result.stderr):
                    copy_tool = "conduit"
                elif shutil.which("pfcp"):
                    copy_tool = "pfcp"
                else:
                    st.error("Data transfer not supported on this HPC system. Ensure you are on a transfer node, not head node.")
                    st.stop()

                # delete data in t2 locations table on scratch before moving to campaign
                store = get_db(temp_t2_db_name)
                df = store.get_table(TIER_2_TABLE, True, True)
                df.iloc[0, 1:] = None # delete all col data except the dsi_table_name col
                store.update(df)
                store.close()

                scratch_move_error = None
                with st.spinner(f"Moving data to HPC campaign with `{copy_tool}`"):
                    try:
                        # TODO: Turn verbose off after done testing
                        s = Sync(temp_t2_db_name, isVerbose=True, skip_index=skip_index, add_dbs=[t1_db_name])
                        s.index(hpc_staging, hpc_campaign)
                        s.copy(copy_tool)
                    except Exception as e:
                        scratch_move_error = e
                
                if scratch_move_error is not None:
                    st.error("Scratch to Campaign Move Error:")
                    if "conduit get" in str(scratch_move_error) and "no credentials" in str(scratch_move_error).lower():
                        st.code("Please run 'conduit get' in the other terminal to be able to move data on this HPC.")
                    else:
                        st.code(str(scratch_move_error))
                    st.stop()

                # register project in diana endpoint
                Path(f"{proj_name}_endpoint.txt").write_text(fed_line, encoding="utf-8")
                Path(f"{proj_name}_endpoint.txt").chmod(0o644)
                if copy_tool == "conduit":
                    result = subprocess.run(["bash", "-lc", "type conduit"], capture_output=True, text=True)
                    copy_command = str(result.stdout).split()
                    for idx, s in enumerate(copy_command):
                        if "/" in s:
                            copy_command = copy_command[idx:idx+3]
                            break
                    copy_command.extend(["cp", f"{proj_name}_endpoint.txt", os.path.join(diana_endpoint, f"{proj_name}_endpoint.txt")])
                else:
                    copy_command = ["pfcp", f"{proj_name}_endpoint.txt", os.path.join(diana_endpoint, f"{proj_name}_endpoint.txt")]
                process = subprocess.Popen(copy_command, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, encoding='latin-1')
                stdout, stderr = process.communicate()
                if process.returncode != 0:
                    st.error("Registering this project in the DIANA catalog failed")
                    st.code(stderr)
                    st.stop()

            else:
                if not st.session_state.local_to_staging_moved:
                    local_move_error = None
                    with st.spinner("Moving local data to HPC staging space - check the terminal for 3 password prompts..."):
                        print(" \n \nMoving local data to HPC staging space - 3 password prompts expected:\n ")
                        try:
                            # TODO: Turn verbose off after done testing
                            s = Sync(temp_t2_db_name, isVerbose=True, skip_index=skip_index, add_dbs=[t1_db_name])
                            s.index(local_data, f"{username}@{hpc_name}:{hpc_staging}")
                            s.copy("rsync")
                        except Exception as e:
                            local_move_error = e
                            print(" \nGo back to app")
                    
                    if local_move_error is not None:
                        st.error("Local Data to HPC Staging Move Error:")
                        st.code(str(local_move_error))
                        st.stop()

                    st.session_state.local_to_staging_moved = True

                if not st.session_state.staging_to_campaign_moved:
                    # run remote script to move from scratch to campaign
                    script = REMOTE_MOVE_SCRIPT
                    full_staging_path = os.path.join(hpc_staging, proj_name)
                    script = script.replace('00000', repr(full_staging_path)) # staging folder
                    script = script.replace('11111', repr(t1_db_name)) # t1 db name
                    script = script.replace('22222', repr(temp_t2_db_name)) # t2 db name
                    script = script.replace('33333', repr(str(Path(local_data).name))) # just name of data folder
                    script = script.replace('44444', repr(hpc_campaign)) # campaign folder
                    script = script.replace('55555', repr(TIER_2_TABLE)) # t2 locations table name

                    with st.spinner("Moving data from HPC staging to HPC campaign — check the terminal for 1 password prompt..."):
                        print(" \n \nMoving data from HPC staging to HPC campaign - 1 password prompt expected:")
                        cmd = ["ssh", f"{username}@{hpc_name}", "python3", "-"]
                        remote_run = subprocess.run(cmd, input=script, text=True, capture_output=True, check=False)
                    print(" \nGo back to app")
                    if remote_run.returncode != 0:
                        st.error("HPC Staging to Campaign Move Error:")
                        if remote_run.stderr:
                            output = remote_run.stdout.lower()
                            if "Add the DSI HPC module to your shell config file" in remote_run.stdout:
                                st.code("Add the DSI HPC module to your shell config file")
                            elif "Only testing on this HPC for now" in remote_run.stdout:
                                st.code("Data transfer not supported on this HPC system. Ensure you are on a transfer node, not head node.")
                            elif "conduit get" in output and ("no credentials" in output or "failed" in output):
                                st.code(f"In a new terminal session on '{hpc_name}', run 'conduit get' to enable data transfer to campaign.")
                            else:
                                st.code(remote_run.stderr)
                        st.stop()

                    if "Only testing on this HPC for now" in remote_run.stdout:
                        st.error("Data transfer not supported on this HPC system. Ensure you are on a transfer node, not head node.")
                        st.stop()
                    elif "Add the DSI HPC module to your shell config file" in remote_run.stdout:
                        st.error("Add the DSI HPC module to your shell config file")
                        st.stop()

                    marker = "DSI HPC Staging to Campaign Move Error"
                    idx = remote_run.stdout.find(marker)
                    if idx != -1:
                        st.error("HPC Staging to Campaign Move Error")
                        output = remote_run.stdout.lower()
                        if "conduit get" in output and ("no credentials" in output or "failed" in output):
                            st.code(f"In a new terminal session on '{hpc_name}', run 'conduit get' to enable data transfer to campaign.")
                        else:
                            st.code(remote_run.stdout[idx + len(marker):])
                        st.stop()
                    
                    st.session_state.staging_to_campaign_moved = True

                endpoint_script = REMOTE_REGISTER_ENDPOINT_SCRIPT
                full_staging_path = os.path.join(hpc_staging, proj_name)
                endpoint_script = endpoint_script.replace('00000', repr(full_staging_path)) # staging folder
                endpoint_script = endpoint_script.replace('11111', repr(fed_line))
                endpoint_script = endpoint_script.replace('22222', repr(proj_name))
                endpoint_script = endpoint_script.replace('33333', repr(diana_endpoint))

                with st.spinner("Registering project in DIANA Catalog — check the terminal for 1 password prompt..."):
                    print(" \n \nRegistering project in DIANA Catalog - 1 password prompt expected:")
                    cmd = ["ssh", f"{username}@{hpc_name}", "python3", "-"]
                    remote_endpoint_run = subprocess.run(cmd, input=endpoint_script, text=True, capture_output=True, check=False)

                if remote_endpoint_run.returncode != 0:
                    st.error("Error registering project in DIANA Catalog:")
                    if remote_endpoint_run.stderr:
                        output = remote_endpoint_run.stdout
                        if "Only testing on this HPC for now" in remote_endpoint_run.stdout:
                            st.code("Cannot register this project on this HPC system. Ensure you are on a transfer node, not head node.")
                        elif "Endpoint error" in remote_endpoint_run.stdout:
                            st.code(remote_endpoint_run.stdout.split("Endpoint error", 1)[1])
                    st.stop()                

            # set user group and access permissions after move
            full_campaign_path = os.path.join(hpc_campaign, proj_name)
            if local_data.lower() == "n/a":
                if user_group:
                    with st.spinner("Updating user group on Campaign"):
                        cmd = ["chgrp", "-R", user_group, full_campaign_path]
                        result = subprocess.run(cmd, capture_output=True, text=True, check=False)
                    if result.returncode != 0:
                        st.error("Error setting user group for data on Campaign")
                        st.code(result.stderr)
                        st.stop()

                if access_permissions_code:
                    with st.spinner("Updating data access permissions on Campaign"):
                        cmd = ["chmod", "-R", access_permissions_code, full_campaign_path]
                        result = subprocess.run(cmd, capture_output=True, text=True, check=False)
                    if result.returncode != 0:
                        st.error("Error setting access permissions for data on Campaign")
                        st.code(result.stderr)
                        st.stop()
            else:
                if user_group and access_permissions_code:
                    print(" \nPassword prompt 1/1: updating user group and data access permissions on HPC campaign")
                    with st.spinner("Updating data group and access permissions on HPC Campaign — check the terminal for 1 password prompt..."):
                        cmd = ["ssh", f"{username}@{hpc_name}",
                            f"chgrp -R {user_group} {full_campaign_path} && chmod -R {access_permissions_code} {full_campaign_path}"]
                        result = subprocess.run(cmd, capture_output=True, text=True, check=False)
                    print(" \nGo back to app")
                    if result.returncode != 0:
                        st.error("Error updating user group and data access permissions on HPC campaign")
                        st.code(result.stderr)
                        st.stop()

            # after move, delete the temp t1 & t2 dbs that were actually moved
            os.remove(temp_t2_db_name)

            # after successful move, set 'has_moved' col for this project in maven.db to be path to local data
            master_store = get_db(get_master_db_name())
            master_df = master_store.find(f"project_id = {qid}", True, True)
            master_df["has_moved"] = local_data if local_data.lower() != "n/a" else hpc_staging
            master_store.update(master_df)
            master_store.close()
            st.session_state.unchanged_data = False
            st.session_state.ran_change_dialog = False
            st.session_state.local_to_staging_moved = False
            st.session_state.staging_to_campaign_moved = False

            st.success(f"Successfully moved data with DSI to HPC campaign: {os.path.join(hpc_campaign, proj_name)}/")

# -----------------------------
# Sidebar
# -----------------------------
st.sidebar.title("Actions")

if st.sidebar.button("➕ New project"):
    # Save current progress
    if st.session_state.active_qid is not None:
        qid = int(st.session_state.active_qid)
        token = _qid_token(qid)
        sec = int(st.session_state.section_idx)
        if sec in ALL_SECTION_IDXS:
            updates = apply_section_updates(sec, token)
            update_datasheet(qid, updates)
        clear_context_file_state(token)
    clear_context_file_state("draft")

    st.session_state.active_qid = None
    st.session_state.section_idx = 0
    st.session_state.draft_answers = {}
    st.session_state.draft_mode = True
    st.session_state.screen = "datasheet"
    st.session_state.ran_change_dialog = False
    st.session_state.local_to_staging_moved = False
    st.session_state.staging_to_campaign_moved = False
    st.session_state.confirm_submit_context_files = False
    st.rerun()

# Button only visible if there is at least one committed project
_committed = list_projects()
if not _committed.empty:
    if st.sidebar.button("📂 Select existing project"):
        # Save current progress
        if st.session_state.active_qid is not None:
            qid = int(st.session_state.active_qid)
            token = _qid_token(qid)
            sec = int(st.session_state.section_idx)
            if sec in ALL_SECTION_IDXS:
                updates = apply_section_updates(sec, token)
                update_datasheet(qid, updates)

        st.session_state.active_qid = None
        st.session_state.section_idx = 0
        st.session_state.draft_answers = {}
        st.session_state.draft_mode = False
        st.session_state.screen = "datasheet"
        st.session_state.ran_change_dialog = False
        st.session_state.local_to_staging_moved = False
        st.session_state.staging_to_campaign_moved = False
        st.session_state.confirm_submit_context_files = False
        st.rerun()

st.sidebar.divider()
st.sidebar.write("### Project Setup")

if st.session_state.active_qid is None:
    label = f"{SECTION_BY_IDX[0]['title']}"
    st.sidebar.button(label, disabled=False,
                      key="nav_draft_0", width="stretch")

    st.sidebar.space()
    st.sidebar.write("### Datasheet Sections")

    for idx in ALL_SECTION_IDXS:
        if idx == 0:
            continue
        label = f"{idx}: {SECTION_BY_IDX[idx]['title']}"
        st.sidebar.button(label, disabled=True,
                          key=f"nav_draft_{idx}", width="stretch")

else:
    qid = int(st.session_state.active_qid)
    token = _qid_token(qid)
    row_latest = get_datasheet(qid)
    max_unlocked = accessed_section_idx(row_latest, "max")

    # Section 0 separated
    is_complete, _ = section_complete(0, row_latest)
    label_prefix = "✅ " if is_complete else "⬜ "
    label = f"{label_prefix} {SECTION_BY_IDX[0]['title']}"

    if st.sidebar.button(label, key=f"nav_{qid}_0", width="stretch"):
        st.session_state.screen = "datasheet"
        if st.session_state.section_idx in ALL_SECTION_IDXS:
            updates = apply_section_updates(
                st.session_state.section_idx, token)
            update_datasheet(qid, updates)
        st.session_state.section_idx = 0
        st.session_state._scroll_to_top = True
        st.session_state.ran_change_dialog = False
        st.session_state.local_to_staging_moved = False
        st.session_state.staging_to_campaign_moved = False
        st.session_state.confirm_submit_context_files = False
        st.rerun()

    st.sidebar.space()
    st.sidebar.write("### Datasheet Sections")

    with st.sidebar:
        with st.container(key="section_btns"):
            for idx in ALL_SECTION_IDXS:
                if idx == 0:
                    continue

                is_complete, _ = section_complete(idx, row_latest)
                label_prefix = "✅ " if is_complete else "⬜ "
                label = f"{label_prefix}{idx}: {SECTION_BY_IDX[idx]['title']}"
                disabled = idx > max_unlocked
                if st.button(label, key=f"nav_{qid}_{idx}", disabled=disabled, width="stretch"):
                    st.session_state.screen = "datasheet"
                    if st.session_state.section_idx in ALL_SECTION_IDXS:
                        updates = apply_section_updates(
                            st.session_state.section_idx, token)
                        update_datasheet(qid, updates)
                    st.session_state.section_idx = idx
                    st.session_state._scroll_to_top = True
                    st.session_state.ran_change_dialog = False
                    st.session_state.local_to_staging_moved = False
                    st.session_state.staging_to_campaign_moved = False
                    st.session_state.confirm_submit_context_files = False
                    st.rerun()

st.sidebar.space()
st.sidebar.write("### Metadata Levels")

is_metadata_btn_disabled = False

if st.session_state.active_qid is None:
    is_metadata_btn_disabled = True
else:
    all_data = get_datasheet(qid)
    if any(not section_complete(idx, all_data)[0] for idx in ALL_SECTION_IDXS):
        is_metadata_btn_disabled = True
    else:
        qid = int(st.session_state.active_qid)
        if not get_tier1_table(qid, check_exists=True):
            is_metadata_btn_disabled = True

if st.session_state.active_qid is None:
    t2_exists = False
else:
    t2_exists = Path(get_tier2_db_path(qid)).is_file()

t1_btn_prefix = "" if is_metadata_btn_disabled or not t2_exists else "✅ "
if st.sidebar.button(f"{t1_btn_prefix}Findability Metadata", disabled=is_metadata_btn_disabled, width="stretch",
                     help="First fill out datasheet sections" if is_metadata_btn_disabled else ""):
    st.session_state.screen = "tier1"
    st.session_state.section_idx = 0
    st.session_state._scroll_to_top = True
    st.session_state.ran_change_dialog = False
    st.session_state.local_to_staging_moved = False
    st.session_state.staging_to_campaign_moved = False
    st.session_state.confirm_submit_context_files = False
    st.rerun()


if st.sidebar.button("AI-Ready Metadata", disabled=is_metadata_btn_disabled or not t2_exists, width="stretch",
                     help="First fill out datasheet sections and tier 1 metadata" if is_metadata_btn_disabled else ""):
    st.session_state.screen = "tier2"
    st.session_state.section_idx = 0
    st.session_state._scroll_to_top = True
    st.session_state.render_t2_extraction = False
    st.session_state.local_to_staging_moved = False
    st.session_state.staging_to_campaign_moved = False
    st.session_state.confirm_submit_context_files = False
    st.rerun()
